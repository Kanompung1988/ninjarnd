"""
NINJA Backend API Server
========================
FastAPI server to connect Python backend with Next.js frontend
Integrated with:
  - deep_research_engine_integrated.py (Unified research from ninja_deepresearch + research_core)
  - unified_presentation_final.py (Slide generation)
"""
import os
import asyncio
from dotenv import load_dotenv
from pathlib import Path

# Load environment variables from .env file
env_path = Path(__file__).parent / '.env'
load_dotenv(dotenv_path=env_path, override=True)

from fastapi import FastAPI, HTTPException, Depends, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
import json
from datetime import datetime
import traceback

# Azure Blob Storage
try:
    from azure.storage.blob import BlobServiceClient, ContentSettings
    AZURE_STORAGE_CONNECTION_STRING = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
    AZURE_STORAGE_AVAILABLE = bool(AZURE_STORAGE_CONNECTION_STRING)
    if AZURE_STORAGE_AVAILABLE:
        blob_service_client = BlobServiceClient.from_connection_string(AZURE_STORAGE_CONNECTION_STRING)
        print("✅ Azure Blob Storage configured")
    else:
        blob_service_client = None
        print("⚠️  Azure Blob Storage not configured - using local storage")
except ImportError:
    print("⚠️  azure-storage-blob not installed - using local storage")
    AZURE_STORAGE_AVAILABLE = False
    blob_service_client = None
except Exception as e:
    print(f"⚠️  Azure Blob Storage initialization failed: {e}")
    AZURE_STORAGE_AVAILABLE = False
    blob_service_client = None

# Import Azure OpenAI Core
try:
    from Azure_OpenAi_core import AzureOpenAICore
    AZURE_OPENAI_AVAILABLE = True
    print("✅ Azure OpenAI Core loaded successfully")
except ImportError as e:
    print(f"⚠️  Azure OpenAI Core import failed: {e}")
    AzureOpenAICore = None
    AZURE_OPENAI_AVAILABLE = False

# Import GLM-4.7 7-Step Generator (Z.AI Direct API) - Primary Slide Generator
try:
    from engines.glm_7step_generator import GLM7StepGenerator
    GLM_7STEP_AVAILABLE = True
    print("✅ GLM-4.7 7-Step Generator loaded (Z.AI API)")
except ImportError as e:
    print(f"⚠️  GLM-4.7 7-Step Generator import failed: {e}")
    GLM7StepGenerator = None
    GLM_7STEP_AVAILABLE = False

# Import Z.AI Fast Slide Generator (Single API call - FAST!)
try:
    from engines.zai_slide_generator import ZAISlideGenerator, generate_presentation_fast
    ZAI_FAST_AVAILABLE = True
    print("✅ Z.AI Fast Slide Generator loaded (Single API)")
except ImportError as e:
    print(f"⚠️  Z.AI Fast Generator import failed: {e}")
    ZAISlideGenerator = None
    generate_presentation_fast = None
    ZAI_FAST_AVAILABLE = False

# Import AI Slide Generator (Z.AI Style with PPTX Export)
try:
    from engines.ai_slide_generator import AISlideGenerator, generate_ai_slides
    AI_SLIDE_AVAILABLE = True
    print("✅ AI Slide Generator loaded (Z.AI Style + PPTX)")
except ImportError as e:
    print(f"⚠️  AI Slide Generator import failed: {e}")
    AISlideGenerator = None
    generate_ai_slides = None
    AI_SLIDE_AVAILABLE = False

# Import comprehensive research engine (unified for all models)
try:
    from engines.comprehensive_research_engine import ComprehensiveResearch, deep_research
    RESEARCH_ENGINE_AVAILABLE = True
    print("✅ comprehensive_research_engine loaded successfully")
except Exception as e:
    print(f"⚠️  comprehensive_research_engine import failed: {e}")
    ComprehensiveResearch = None
    deep_research = None
    RESEARCH_ENGINE_AVAILABLE = False

# Presentation export using python-pptx directly
PRESENTATION_ENGINE_AVAILABLE = True

# Export function - moved from presentation_engine.py
def export_presentation(presentation_data: Dict[str, Any], format: str = "pptx") -> Path:
    """Export presentation to PPTX format"""
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    if format.lower() != "pptx":
        raise ValueError(f"Unsupported format: {format}. Use 'pptx'.")
    
    exports_dir = Path(__file__).parent / "exports" / "presentations"
    exports_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        prs = PptxPresentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_layout)
        title = presentation_data.get('title', 'Presentation')
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Content slides
        for slide_data in presentation_data.get('slides', []):
            layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(layout)
            
            # Title
            slide_title = slide_data.get('title', '')
            if slide_title:
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(36)
                p.font.bold = True
            
            # Content
            content = slide_data.get('content', '')
            if content:
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(18)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}.pptx"
        filepath = exports_dir / filename
        prs.save(str(filepath))
        print(f"✅ PPTX exported: {filepath}")
        return filepath
        
    except Exception as e:
        print(f"❌ PPTX export error: {e}")
        raise

app = FastAPI(title="NINJA Backend API", version="1.0.0")

# CORS configuration for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:3001",
        "http://localhost:3002",
        "http://127.0.0.1:3000",
        "http://127.0.0.1:3001",
        "http://127.0.0.1:3002",
        "https://*.vercel.app",
        "*"  # Allow all origins for development
    ],
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS", "PATCH"],
    allow_headers=["*"],
    expose_headers=["*"],
)

# API Key verification (optional - ทำให้ flexible สำหรับ frontend)
BACKEND_API_SECRET = os.getenv("BACKEND_API_SECRET", "dev-secret-key")

def verify_api_key_optional(authorization: str = Header(None)):
    """Optional API key verification - ไม่ raise error ถ้าไม่มี"""
    if authorization and authorization.startswith("Bearer "):
        token = authorization.replace("Bearer ", "")
        if token == BACKEND_API_SECRET:
            return token
    return None

def verify_api_key_required(authorization: str = Header(None)):
    """Required API key verification - raise error ถ้าไม่มีหรือผิด"""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing authorization header")
    token = authorization.replace("Bearer ", "")
    if token != BACKEND_API_SECRET:
        raise HTTPException(status_code=401, detail="Invalid API key")
    return token

# Storage directories
USER_DATA_DIR = Path("user_data")
USER_DATA_DIR.mkdir(exist_ok=True)
RESEARCH_BLOGS_DIR = Path("research_blogs")
PRESENTATIONS_DIR = Path("presentations")
RESEARCH_BLOGS_DIR.mkdir(exist_ok=True)
PRESENTATIONS_DIR.mkdir(exist_ok=True)

# Helper functions for user-specific storage
def get_user_folder(user_id: str, data_type: str) -> Path:
    """Get user-specific folder path"""
    # Sanitize user_id for filesystem
    safe_user_id = user_id.replace("@", "_at_").replace(".", "_")
    user_folder = USER_DATA_DIR / safe_user_id / data_type
    user_folder.mkdir(parents=True, exist_ok=True)
    return user_folder

def save_user_data_to_blob(user_id: str, data_type: str, file_name: str, data: bytes):
    """Save user data to Azure Blob Storage"""
    if not AZURE_STORAGE_AVAILABLE or not blob_service_client:
        return None
    
    try:
        container_name = f"user-{data_type}"
        blob_name = f"{user_id}/{file_name}"
        
        # Create container if not exists
        try:
            blob_service_client.create_container(container_name)
        except:
            pass  # Container already exists
        
        # Upload blob
        blob_client = blob_service_client.get_blob_client(
            container=container_name,
            blob=blob_name
        )
        blob_client.upload_blob(data, overwrite=True)
        return blob_client.url
    except Exception as e:
        print(f"Error saving to blob: {e}")
        return None

def load_user_data_from_blob(user_id: str, data_type: str, file_name: str) -> Optional[bytes]:
    """Load user data from Azure Blob Storage"""
    if not AZURE_STORAGE_AVAILABLE or not blob_service_client:
        return None
    
    try:
        container_name = f"user-{data_type}"
        blob_name = f"{user_id}/{file_name}"
        
        blob_client = blob_service_client.get_blob_client(
            container=container_name,
            blob=blob_name
        )
        return blob_client.download_blob().readall()
    except Exception as e:
        print(f"Error loading from blob: {e}")
        return None

# Helper function to check if model is Azure OpenAI
def is_azure_model(model: str) -> bool:
    """Check if the model should use Azure OpenAI"""
    azure_prefixes = [
        "gpt-5", "gpt-5-", 
        "o1", "o1-",
        "gpt-4o", "gpt-4o-"
    ]
    return any(model.startswith(prefix) for prefix in azure_prefixes)

class ChatRequest(BaseModel):
    message: str
    chat_history: List[Dict[str, Any]] = []
    chat_id: str = "default"
    user_id: Optional[str] = "default@user.com"  # User email (optional)
    user_name: Optional[str] = None
    deep_research_mode: bool = False
    realtime_research_mode: bool = False
    agent_mode: bool = False
    model: str = "gemini-2.0-flash-exp"  # Changed from typhoon to gemini (more reliable)
    search_engine: str = "hybrid"

class ResearchRequest(BaseModel):
    topic: str
    user_id: Optional[str] = "default@user.com"  # User email (optional)
    user_name: Optional[str] = None
    days_back: int = 7
    effort: str = "standard"
    scope: str = "balanced"
    model: str = "gemini-2.0-flash-exp"  # Changed from typhoon to gemini (more reliable)
    chat_id: str = "default"
    search_engine: str = "hybrid"
    use_hybrid_search: bool = True

class SlideRequest(BaseModel):
    chat_id: str
    user_id: Optional[str] = "default@user.com"  # User email (optional)
    user_name: Optional[str] = None
    json_path: str

class PresentationGenerateRequest(BaseModel):
    topic: str
    user_id: Optional[str] = "default@user.com"  # User email (optional)
    user_name: Optional[str] = None
    slide_count: int = 8
    style: str = "professional"
    aspect_ratio: str = "16:9"
    model: str = "gemini-2.0-flash-exp"
    fontSizeMode: str = "optimized"
    researchContext: Optional[str] = None
    generate_images: bool = True

class PresentationSaveRequest(BaseModel):
    title: str
    user_id: Optional[str] = "default@user.com"  # User email (optional)
    user_name: Optional[str] = None
    slides: List[Dict[str, Any]]
    metadata: Dict[str, Any] = {}

class ResearchBlogSaveRequest(BaseModel):
    query: str
    user_id: Optional[str] = "default@user.com"  # User email (optional)
    user_name: Optional[str] = None
    content: str
    sources: List[Dict[str, Any]] = []
    metadata: Dict[str, Any] = {}
    model: str = "gemini-2.0-flash-exp"

class PresentationExportRequest(BaseModel):
    presentation: Dict[str, Any]
    user_id: Optional[str] = "default@user.com"  # User email (optional)
    user_name: Optional[str] = None
    format: str = "pptx"
    output_format: str = "both"
    settings: Dict[str, Any] = {}

class ResearchBlogSaveRequest(BaseModel):
    title: str
    query: str
    content: str
    sources: List[Dict[str, Any]] = []
    metadata: Dict[str, Any] = {}

# Directory for storing research blogs
RESEARCH_BLOGS_DIR = Path(__file__).parent / "research_blogs"
RESEARCH_BLOGS_DIR.mkdir(parents=True, exist_ok=True)

# Endpoints
@app.get("/")
async def root():
    return {
        "message": "NINJA Backend API - Integrated with DeepResearch",
        "version": "1.0.0",
        "status": "running",
        "features": ["chat", "deep_research", "presentations"]
    }

async def simple_chat(request: ChatRequest):
    """Simple chat mode when research engine is not available"""
    try:
        # NINJA System Prompt
        ninja_system = """
You are **NINJA**, an intelligent multi-agent assistant designed to transform diverse data sources into actionable insights.

🥷 **NINJA CAPABILITIES**
- Multi-agent LLM framework for information processing  
- Transform news, research papers, social media, code, and databases into concise briefs  
- Topic-guided prompts with ReACT-style reasoning  
- Extract key facts, filter noise, and validate accuracy  
- Automate the end-to-end process from search to presentation  
- Reduce information overload and accelerate decision-making  
- Adapt seamlessly to competitive intelligence, policy monitoring, education, and R&D  

🎯 **NINJA PERSONALITY**
- Professional yet approachable  
- Concise and action-oriented  
- Expertise in data analysis and strategic insights  
- Always provide structured, slide-ready information when possible  

💡 **HOW TO SWITCH MODES**
- To enable **DeepResearch Mode**, click the **"+" Toggle button** on the right side of the chat input to switch modes.
"""
        
        # Model name mapping for backward compatibility
        model_mapping = {
            "typhoon-v2.5-instruct": "typhoon-v2.5-30b-a3b-instruct",
            "typhoon-v2-instruct": "typhoon-v2.5-30b-a3b-instruct",
        }
        
        # Map old model names to new ones
        actual_model = model_mapping.get(request.model, request.model)
        
        # Check if Ensemble model
        if actual_model.startswith("ensemble-"):
            print(f"🔀 Using Ensemble Model: {actual_model}")
            ensemble_answers = []
            
            if actual_model == "ensemble-thai-expert":
                # Typhoon 30B + Gemini Flash
                print("  📊 Running Typhoon 30B...")
                try:
                    from openai import OpenAI
                    import httpx
                    http_client = httpx.Client(timeout=httpx.Timeout(10.0, connect=5.0))
                    typhoon_client = OpenAI(
                        api_key=os.getenv("TYPHOON_API_KEY"),
                        base_url="https://api.opentyphoon.ai/v1",
                        http_client=http_client
                    )
                    typhoon_response = typhoon_client.chat.completions.create(
                        model="typhoon-v2.5-30b-a3b-instruct",
                        messages=[{"role": "system", "content": ninja_system}, {"role": "user", "content": request.message}],
                        timeout=30.0
                    )
                    ensemble_answers.append(("Typhoon 30B", typhoon_response.choices[0].message.content))
                except Exception as e:
                    print(f"  ⚠️ Typhoon error (skipping): {e}")
                
                print("  📊 Running Gemini Flash...")
                try:
                    import google.generativeai as genai
                    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
                    gemini_model = genai.GenerativeModel("gemini-2.0-flash-exp")
                    gemini_response = gemini_model.generate_content(f"{ninja_system}\n\nUser: {request.message}")
                    ensemble_answers.append(("Gemini Flash", gemini_response.text))
                except Exception as e:
                    print(f"  ⚠️ Gemini error: {e}")
            
            elif actual_model == "ensemble-balanced":
                # Typhoon 30B + GPT-4o
                print("  📊 Running Typhoon 30B...")
                try:
                    from openai import OpenAI
                    import httpx
                    http_client = httpx.Client(timeout=httpx.Timeout(10.0, connect=5.0))
                    typhoon_client = OpenAI(
                        api_key=os.getenv("TYPHOON_API_KEY"),
                        base_url="https://api.opentyphoon.ai/v1",
                        http_client=http_client
                    )
                    typhoon_response = typhoon_client.chat.completions.create(
                        model="typhoon-v2.5-30b-a3b-instruct",
                        messages=[{"role": "system", "content": ninja_system}, {"role": "user", "content": request.message}],
                        timeout=30.0
                    )
                    ensemble_answers.append(("Typhoon 30B", typhoon_response.choices[0].message.content))
                except Exception as e:
                    print(f"  ⚠️ Typhoon error (skipping): {e}")
                
                print("  📊 Running GPT-4o...")
                try:
                    azure_client = AzureOpenAICore()
                    gpt4o_response = azure_client.chat(
                        messages=[{"role": "system", "content": ninja_system}, {"role": "user", "content": request.message}],
                        model="gpt-4o"
                    )
                    ensemble_answers.append(("GPT-4o", gpt4o_response['content']))
                except Exception as e:
                    print(f"  ⚠️ GPT-4o error: {e}")
            
            elif actual_model == "ensemble-deep-reasoning":
                # Typhoon 30B + O1
                print("  📊 Running Typhoon 30B...")
                try:
                    from openai import OpenAI
                    import httpx
                    http_client = httpx.Client(timeout=httpx.Timeout(10.0, connect=5.0))
                    typhoon_client = OpenAI(
                        api_key=os.getenv("TYPHOON_API_KEY"),
                        base_url="https://api.opentyphoon.ai/v1",
                        http_client=http_client
                    )
                    typhoon_response = typhoon_client.chat.completions.create(
                        model="typhoon-v2.5-30b-a3b-instruct",
                        messages=[{"role": "system", "content": ninja_system}, {"role": "user", "content": request.message}],
                        timeout=30.0
                    )
                    ensemble_answers.append(("Typhoon 30B", typhoon_response.choices[0].message.content))
                except Exception as e:
                    print(f"  ⚠️ Typhoon error (skipping): {e}")
                
                print("  📊 Running O1...")
                try:
                    azure_client = AzureOpenAICore()
                    o1_response = azure_client.chat(
                        messages=[{"role": "user", "content": request.message}],
                        model="o1"
                    )
                    ensemble_answers.append(("O1", o1_response['content']))
                except Exception as e:
                    print(f"  ⚠️ O1 error: {e}")
            
            # Synthesize ensemble responses
            if ensemble_answers:
                synthesis_prompt = f"""You are an expert synthesizer. Multiple AI models have answered the same question. 
Combine their insights into ONE comprehensive, accurate answer.

Question: {request.message}

Answers from models:
"""
                for model_name, ans in ensemble_answers:
                    synthesis_prompt += f"\n### {model_name}:\n{ans}\n"
                
                synthesis_prompt += "\nProvide a synthesized answer that combines the best insights from all models:"
                
                # Use GPT-4o for synthesis
                try:
                    azure_client = AzureOpenAICore()
                    final_response = azure_client.chat(
                        messages=[{"role": "user", "content": synthesis_prompt}],
                        model="gpt-4o"
                    )
                    answer = final_response['content']
                except Exception as e:
                    print(f"  ⚠️ Synthesis error: {e}")
                    # Fallback: just concatenate answers
                    answer = "\n\n---\n\n".join([f"**{name}:**\n{ans}" for name, ans in ensemble_answers])
            else:
                answer = "Ensemble failed - no models responded successfully"
        
        # Check if Azure OpenAI model
        elif is_azure_model(actual_model) and AZURE_OPENAI_AVAILABLE:
            print(f"🔵 Using Azure OpenAI: {actual_model}")
            azure_client = AzureOpenAICore()
            response = azure_client.chat(
                messages=[
                    {"role": "system", "content": ninja_system},
                    {"role": "user", "content": request.message}
                ],
                model=actual_model
            )
            answer = response['content']
            
        elif actual_model.startswith("typhoon"):
            print(f"⚠️  Attempting Typhoon API with timeout protection...")
            try:
                from openai import OpenAI
                import httpx
                # Add aggressive timeout: 10 seconds for connection, 30 seconds total
                http_client = httpx.Client(timeout=httpx.Timeout(10.0, connect=5.0))
                client = OpenAI(
                    api_key=os.getenv("TYPHOON_API_KEY"),
                    base_url="https://api.opentyphoon.ai/v1",
                    http_client=http_client
                )
                response = client.chat.completions.create(
                    model=actual_model,
                    messages=[
                        {"role": "system", "content": ninja_system},
                        {"role": "user", "content": request.message}
                    ],
                    timeout=30.0  # Additional timeout parameter
                )
                answer = response.choices[0].message.content
            except Exception as typhoon_error:
                print(f"❌ Typhoon API failed: {typhoon_error}")
                print(f"🔄 Falling back to Gemini...")
                # Fallback to Gemini
                import google.generativeai as genai
                genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
                model = genai.GenerativeModel("gemini-2.0-flash-exp")
                enhanced_message = f"{ninja_system}\n\nUser: {request.message}"
                response = model.generate_content(enhanced_message)
                answer = f"⚠️ (Answered by Gemini - Typhoon unavailable)\n\n{response.text}"
            
        elif request.model.startswith("gpt"):
            from openai import OpenAI
            client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            response = client.chat.completions.create(
                model=request.model,
                messages=[
                    {"role": "system", "content": ninja_system},
                    {"role": "user", "content": request.message}
                ]
            )
            answer = response.choices[0].message.content
            
        elif request.model.startswith("gemini"):
            import google.generativeai as genai
            genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
            model = genai.GenerativeModel(request.model)
            # Gemini doesn't use system prompt in the same way, so prepend to message
            enhanced_message = f"{ninja_system}\n\nUser: {request.message}"
            response = model.generate_content(enhanced_message)
            answer = response.text
            
        else:
            answer = f"Model {request.model} not supported in simple chat mode"
        
        return {
            "success": True,
            "mode": "simple_chat",
            "query": request.message,
            "model": request.model,
            "response": answer,
            "brief": {
                "direct_answer": answer[:500],
                "key_findings": [],
                "confidence": "medium",
                "recommendations": []
            },
            "sources": [],
            "credibility_assessment": {},
            "metadata": {"mode": "simple_chat"},
            "total_sources": 0,
            "pipeline_stages": ["Direct LLM Response"]
        }
    except Exception as e:
        print(f"Simple chat error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/chat")
async def chat(
    request: ChatRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """Chat with NINJA AI - Multi-Mode Support (Normal, Full DeepResearch, Realtime DeepResearch, Agent)"""
    try:
        # Check if user wants Full DeepResearch mode (with blog generation)
        if request.deep_research_mode and RESEARCH_ENGINE_AVAILABLE:
            print(f"🔬 Full DeepResearch Mode: {request.model}")
            start_time = datetime.now()
            
            # Use comprehensive research engine with extended search
            engine = ComprehensiveResearch(
                model=request.model,
                tavily_api_key=os.getenv("TAVILY_API_KEY"),
                serper_api_key=os.getenv("SERPER_API_KEY"),
                jina_api_key=os.getenv("JINA_API_KEY"),
                openai_api_key=os.getenv("OPENAI_API_KEY"),
                typhoon_api_key=os.getenv("TYPHOON_API_KEY"),
                gemini_api_key=os.getenv("GEMINI_API_KEY")
            )
            
            # Execute comprehensive research pipeline - NO time limit for thorough analysis
            result = engine.research(
                query=request.message,
                recency_days=365,  # Search all time for comprehensive coverage
                export_pptx=False
            )
            
            # Save research content as blog for future slide generation
            # ✅ Use user-specific folder
            user_blog_dir = get_user_folder(request.user_id, "research_blogs")
            
            # Create unique filename based on timestamp and query
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_query = "".join(c for c in request.message[:50] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
            blog_filename = f"{timestamp}_{safe_query}.json"
            blog_path = user_blog_dir / blog_filename
            
            # Save research data
            blog_data = {
                "timestamp": timestamp,
                "query": request.message,
                "model": request.model,
                "markdown_report": result.markdown_report,
                "executive_summary": result.executive_summary,
                "key_findings": result.key_findings,
                "recommendations": result.recommendations,
                "sources": result.sources,
                "credibility_assessment": result.credibility_assessment,
                "metadata": result.metadata,
                "chat_id": request.chat_id
            }
            
            with open(blog_path, 'w', encoding='utf-8') as f:
                json.dump(blog_data, f, indent=2, ensure_ascii=False)
            
            print(f"📝 Saved research blog: {blog_path}")
            
            # Calculate duration and stats
            duration = (datetime.now() - start_time).total_seconds()
            avg_credibility = sum(s.get('credibility_score', 0.5) for s in result.sources) / len(result.sources) if result.sources else 0.5
            
            # Format response with research_data for beautiful display
            return {
                "success": True,
                "mode": "deep_research",
                "query": request.message,
                "model": request.model,
                "response": result.markdown_report,
                "blog_id": blog_filename,  # Add blog ID for reference
                "brief": {
                    "direct_answer": result.executive_summary,
                    "key_findings": result.key_findings,
                    "confidence": result.metadata.get('confidence', 'high'),
                    "recommendations": result.recommendations
                },
                "sources": result.sources,
                "credibility_assessment": result.credibility_assessment,
                "metadata": result.metadata,
                "total_sources": len(result.sources),
                "research_data": {
                    "query": request.message,
                    "mode": "deep",
                    "summary": result.executive_summary,
                    "executive_summary": result.executive_summary,
                    "key_findings": result.key_findings,
                    "detailed_analysis": result.detailed_analysis if hasattr(result, 'detailed_analysis') else "",
                    "recommendations": result.recommendations,
                    "sources": [
                        {
                            "title": src.get('title', 'Untitled'),
                            "url": src.get('url', '#'),
                            "snippet": src.get('snippet', 'No description'),
                            "source": src.get('source', 'Web'),
                            "credibility": src.get('credibility_score', 0.5),
                            "date": src.get('published_date', 'N/A')
                        }
                        for src in result.sources
                    ],
                    "total_sources": len(result.sources),
                    "confidence": result.metadata.get('confidence', 'HIGH'),
                    "timestamp": datetime.now().isoformat(),
                    "model": request.model,
                    "average_credibility": avg_credibility,
                    "search_duration": duration
                },
                "pipeline_stages": [
                    "Query Expansion",
                    "Multi-Search (Tavily + SerperAPI + JINA)",
                    "Noise Filtering",
                    "Credibility Scoring",
                    "Synthesis",
                    "Fact Validation",
                    "Sensitive Data Redaction",
                    "Report Generation"
                ],
                "chat_id": request.chat_id,
                "can_generate_slides": True
            }
        
        # Check if user wants Realtime DeepResearch mode (quick, focused on latest news with 20+ sources)
        elif request.realtime_research_mode and RESEARCH_ENGINE_AVAILABLE:
            print(f"⚡ Realtime DeepResearch Mode: {request.model}")
            start_time = datetime.now()
            
            # Use comprehensive research engine with extended sources
            engine = ComprehensiveResearch(
                model=request.model,
                tavily_api_key=os.getenv("TAVILY_API_KEY"),
                serper_api_key=os.getenv("SERPER_API_KEY"),
                jina_api_key=os.getenv("JINA_API_KEY"),
                openai_api_key=os.getenv("OPENAI_API_KEY"),
                typhoon_api_key=os.getenv("TYPHOON_API_KEY"),
                gemini_api_key=os.getenv("GEMINI_API_KEY")
            )
            
            # Execute research for latest news (3 days to get more sources)
            result = engine.research(
                query=request.message,
                recency_days=3,  # Extend to 3 days for more sources
                export_pptx=False
            )
            
            # Sort sources by date (newest first) and take top 20+
            sorted_sources = sorted(
                result.sources,
                key=lambda x: x.get('published_date') or '1970-01-01',
                reverse=True
            )[:25]  # Limit to 25 sources
            
            # Generate key summary for each source using LLM
            source_summaries = []
            for i, src in enumerate(sorted_sources, 1):
                snippet = src.get('snippet', '') or src.get('content', '')[:500]
                title = src.get('title', 'Untitled')
                pub_date = src.get('published_date', 'N/A')
                credibility = src.get('credibility_score', 0)
                
                # Create structured source summary
                source_summaries.append({
                    "index": i,
                    "title": title,
                    "url": src.get('url', '#'),
                    "date": pub_date,
                    "credibility": credibility,
                    "snippet": snippet[:300],
                    "key_point": f"� {snippet[:200]}..." if len(snippet) > 200 else f"📌 {snippet}",
                    "source_type": src.get('source', 'Web')
                })
            
            # Calculate duration
            duration = (datetime.now() - start_time).total_seconds()
            avg_credibility = sum(s['credibility'] for s in source_summaries) / len(source_summaries) if source_summaries else 0.5
            
            # Build structured response
            return {
                "success": True,
                "mode": "realtime_research",
                "query": request.message,
                "model": request.model,
                "response": result.markdown_report,  # Use full report
                "brief": {
                    "direct_answer": result.executive_summary,
                    "key_findings": result.key_findings[:5],
                    "confidence": result.metadata.get('confidence', 'medium'),
                    "recency_focus": "3_days",
                    "total_sources": len(source_summaries)
                },
                "sources": sorted_sources,
                "total_sources": len(source_summaries),
                "research_data": {
                    "query": request.message,
                    "mode": "realtime",
                    "summary": result.executive_summary,
                    "executive_summary": result.executive_summary,
                    "key_findings": result.key_findings,
                    "detailed_analysis": result.detailed_analysis if hasattr(result, 'detailed_analysis') else "",
                    "recommendations": result.recommendations if hasattr(result, 'recommendations') else [],
                    "sources": [
                        {
                            "title": s["title"],
                            "url": s["url"],
                            "snippet": s["snippet"],
                            "key_point": s["key_point"],
                            "source": s["source_type"],
                            "credibility": s["credibility"],
                            "date": s["date"]
                        }
                        for s in source_summaries
                    ],
                    "total_sources": len(source_summaries),
                    "confidence": result.metadata.get('confidence', 'MEDIUM'),
                    "timestamp": datetime.now().isoformat(),
                    "model": request.model,
                    "search_duration": duration,
                    "average_credibility": avg_credibility,
                    "recency_days": 3
                },
                "metadata": {
                    **result.metadata,
                    "mode": "realtime",
                    "blog_generated": False,
                    "recency_days": 3,
                    "focus": "latest_news_extended",
                    "can_generate_slides": True
                },
                "chat_id": request.chat_id
            }

        
        # Check if user wants Agent mode
        elif request.agent_mode:
            print(f"🤖 Agent Mode: {request.model}")
            # Agent mode placeholder - to be implemented
            return {
                "success": True,
                "mode": "agent_mode",
                "query": request.message,
                "model": request.model,
                "response": "🤖 Agent mode is under development. This will enable autonomous AI actions and decision-making capabilities.",
                "metadata": {
                    "mode": "agent",
                    "status": "coming_soon"
                }
            }
        
        else:
            # Normal chat mode (no research)
            print(f"💬 Normal Chat Mode: {request.model}")
            return await simple_chat(request)
    except Exception as e:
        print(f"Chat error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/research")
async def research(
    request: ResearchRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """Perform comprehensive deep research with NINJA (always uses full pipeline)"""
    try:
        if not RESEARCH_ENGINE_AVAILABLE:
            raise HTTPException(status_code=500, detail="Research engine not available")
        
        print(f"🔬 Comprehensive Research: {request.model}")
        
        # Create research engine
        engine = ComprehensiveResearch(
            model=request.model,
            tavily_api_key=os.getenv("TAVILY_API_KEY"),
            serper_api_key=os.getenv("SERPER_API_KEY"),
            jina_api_key=os.getenv("JINA_API_KEY"),
            openai_api_key=os.getenv("OPENAI_API_KEY"),
            typhoon_api_key=os.getenv("TYPHOON_API_KEY"),
            gemini_api_key=os.getenv("GEMINI_API_KEY")
        )
        
        # Execute comprehensive research
        result = engine.research(
            query=request.topic,
            recency_days=request.days_back,
            export_pptx=True  # Export JSON for slide generation
        )
        
        return {
            "success": True,
            "data": {
                "query": result.query,
                "executive_summary": result.executive_summary,
                "key_findings": result.key_findings,
                "detailed_analysis": result.detailed_analysis,
                "sources": result.sources,
                "credibility_assessment": result.credibility_assessment,
                "recommendations": result.recommendations,
                "markdown_report": result.markdown_report,
                "powerpoint_path": result.powerpoint_path,
                "metadata": result.metadata
            },
            "chat_id": request.chat_id,
            "model": request.model
        }
    except Exception as e:
        print(f"Research error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/slides")
async def generate_slides(
    request: SlideRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """Generate presentation slides from research data using Gemini only"""
    try:
        print("\n" + "="*70)
        print("🎨 SLIDE GENERATION REQUEST")
        print("="*70)
        print(f"Chat ID: {request.chat_id}")
        print(f"JSON Path: {request.json_path}")
        print(f"Output Format: {request.output_format}")
        print("="*70 + "\n")
        
        # Use GLM7StepGenerator (primary slide generator)
        if not GLM_7STEP_AVAILABLE:
            raise HTTPException(status_code=503, detail="GLM7StepGenerator not available")
        
        # Load research data
        json_path = Path(request.json_path)
        if not json_path.exists():
            raise HTTPException(status_code=404, detail=f"Research data not found: {request.json_path}")
        
        with open(json_path, 'r', encoding='utf-8') as f:
            research_data = json.load(f)
        
        # Create output directory
        # ✅ Use user-specific folder for chat sessions
        user_chat_dir = get_user_folder(request.user_id, "chat_sessions")
        chat_folder = user_chat_dir / request.chat_id / "outputs"
        chat_folder.mkdir(parents=True, exist_ok=True)
        
        # Generate timestamp for filename
        ts = datetime.now().strftime("%Y%m%dT%H%M%S")
        output_path = chat_folder / f"presentation_{ts}.pptx"
        
        # Initialize GLM7Step generator
        generator = GLM7StepGenerator()
        
        # Extract topic from research data
        topic = research_data.get("query", research_data.get("title", "Research Presentation"))
        context = research_data.get("synthesis", research_data.get("content", ""))
        
        # Generate presentation using GLM
        result = generator.generate_presentation_fast(
            topic=topic,
            slide_count=8,
            style="professional",
            research_context=str(context)[:3000] if context else None
        )
        
        if result.get('success'):
            return {
                "success": True,
                "data": {
                    "output_path": str(output_path),
                    "num_slides": len(result.get('slides', [])),
                    "file_size": 0,
                    "format": "PowerPoint 16:9 (GLM-4.7 Generated)",
                    "generated_at": datetime.now().isoformat(),
                    "download_url": f"/api/slides/download?chat_id={request.chat_id}&file_type=pptx"
                },
                "chat_id": request.chat_id,
                "message": f"✅ Successfully generated {len(result.get('slides', []))} slides using GLM-4.7"
            }
        else:
            raise HTTPException(status_code=500, detail=result.get('error', 'Unknown error'))
            
    except Exception as e:
        print(f"❌ Slides generation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/slides/download")
async def download_slides(
    chat_id: str,
    file_type: str,
    user_id: str,
    api_key: str = Depends(verify_api_key_optional)
):
    """Download generated presentation file"""
    try:
        # Find the file in chat session folder
        # ✅ Use user-specific folder
        user_chat_dir = get_user_folder(user_id, "chat_sessions")
        chat_folder = user_chat_dir / chat_id / "outputs"
        
        if file_type == "pptx":
            files = list(chat_folder.glob("*.pptx"))
            if not files:
                raise HTTPException(status_code=404, detail="PPTX file not found")
            file_path = files[0]
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif file_type == "html":
            files = list(chat_folder.glob("*.html"))
            if not files:
                raise HTTPException(status_code=404, detail="HTML file not found")
            file_path = files[0]
            media_type = "text/html"
        else:
            raise HTTPException(status_code=400, detail="Invalid file type")
        
        return FileResponse(
            path=file_path,
            media_type=media_type,
            filename=f"presentation_{chat_id}.{file_type}"
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"Download error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/health")
async def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "services": {
            "research_engine": RESEARCH_ENGINE_AVAILABLE,
            "presentation_engine": PRESENTATION_ENGINE_AVAILABLE,
        }
    }

@app.post("/api/research/save-blog")
async def save_research_blog(
    request: ResearchBlogSaveRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Save DeepResearch results as a blog for later use in slide generation
    """
    try:
        print(f"💾 Saving research blog: {request.title} for user: {request.user_id}")
        
        # ✅ Use user-specific folder
        user_blog_dir = get_user_folder(request.user_id, "research_blogs")
        
        # Create a unique filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c for c in request.title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        safe_title = safe_title.replace(' ', '_')
        filename = f"{timestamp}_{safe_title}.json"
        filepath = user_blog_dir / filename
        
        # Prepare blog data
        blog_data = {
            "id": timestamp,
            "title": request.title,
            "query": request.query,
            "content": request.content,
            "sources": request.sources,
            "metadata": {
                **request.metadata,
                "saved_at": datetime.now().isoformat(),
                "filename": filename
            }
        }
        
        # Save to file
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(blog_data, f, ensure_ascii=False, indent=2)
        
        print(f"✅ Saved research blog: {filepath}")
        
        return {
            "success": True,
            "blog_id": timestamp,
            "filename": filename,
            "message": "Research blog saved successfully"
        }
        
    except Exception as e:
        print(f"❌ Save blog error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/research/blogs")
async def list_research_blogs(
    user_id: str,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    List all saved research blogs for specific user
    """
    try:
        # ✅ Use user-specific folder
        user_blog_dir = get_user_folder(user_id, "research_blogs")
        blogs = []
        
        for filepath in sorted(user_blog_dir.glob("*.json"), reverse=True):
            try:
                with open(filepath, 'r', encoding='utf-8') as f:
                    blog_data = json.load(f)
                    blogs.append({
                        "id": blog_data.get("id"),
                        "title": blog_data.get("title"),
                        "query": blog_data.get("query"),
                        "saved_at": blog_data.get("metadata", {}).get("saved_at"),
                        "filename": blog_data.get("metadata", {}).get("filename"),
                        "source_count": len(blog_data.get("sources", []))
                    })
            except Exception as e:
                print(f"Error reading blog {filepath}: {e}")
                continue
        
        return {
            "success": True,
            "blogs": blogs,
            "total": len(blogs)
        }
        
    except Exception as e:
        print(f"❌ List blogs error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/research/blogs/{blog_id}")
async def get_research_blog(
    blog_id: str,
    user_id: str,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Get a specific research blog by ID for specific user
    """
    try:
        # ✅ Use user-specific folder
        user_blog_dir = get_user_folder(user_id, "research_blogs")
        
        # Find the blog file
        blog_files = list(user_blog_dir.glob(f"{blog_id}_*.json"))
        
        if not blog_files:
            raise HTTPException(status_code=404, detail="Blog not found")
        
        filepath = blog_files[0]
        
        with open(filepath, 'r', encoding='utf-8') as f:
            blog_data = json.load(f)
        
        return {
            "success": True,
            "blog": blog_data
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Get blog error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/config/models")
async def get_available_models():
    """Get available LLM models and their configuration"""
    return {
        "models": {
            "glm-4.6": {
                "name": "GLM-4.6 (Z.AI)",
                "provider": "Z.AI",
                "available": bool(os.getenv("GLM_API_KEY")),
                "features": ["chat", "reasoning", "coding", "presentations"],
                "context": "200K tokens"
            },
            "gemini-2.0-flash-exp": {
                "name": "Gemini 2.0 Flash",
                "provider": "Google",
                "available": bool(os.getenv("GEMINI_API_KEY"))
            },
            "typhoon-v2.5-instruct": {
                "name": "Typhoon v2.5 Instruct",
                "provider": "OpenTyphoon",
                "available": bool(os.getenv("TYPHOON_API_KEY"))
            },
            "gpt-4": {
                "name": "GPT-4",
                "provider": "OpenAI",
                "available": bool(os.getenv("OPENAI_API_KEY"))
            }
        }
    }

@app.get("/api/models")
async def get_models():
    """Get available LLM models for DeepResearch"""
    return {
        "success": True,
        "models": {
            "glm-4.7": "GLM-4.7 (HuggingFace) - 7-Step Process"
        }
    }

@app.get("/api/slide-generators")
async def get_slide_generators():
    """Get available slide generators - GLM-4.7 7-Step Only"""
    generators = {}
    
    # Check GLM-4.7 7-Step Generator
    if GLM_7STEP_AVAILABLE:
        generators["glm-7step"] = {
            "name": "GLM-4.7 (7-Step)",
            "status": "available",
            "description": "GLM-4.7 via Z.AI Direct API with 7-Step Process",
            "process": [
                "1. Content Analysis",
                "2. Planning & Direction",
                "3. Research & Data",
                "4. Image Generation",
                "5. Create Slides",
                "6. Review & Improve",
                "7. Finalize"
            ]
        }
    else:
        generators["glm-7step"] = {
            "name": "GLM-4.7 (7-Step)",
            "status": "unavailable",
            "description": "GLM_API_KEY not configured"
        }
    
    return {
        "success": True,
        "generators": generators,
        "recommended": "glm-7step" if GLM_7STEP_AVAILABLE else None
    }

@app.get("/api/search-engines")
async def get_search_engines():
    """Get available search engines - Always uses Tavily + SerperAPI + JINA"""
    return {
        "success": True,
        "search_engines": {
            "hybrid": "Hybrid (Tavily + SerperAPI + JINA) - Comprehensive",
        },
        "note": "All models use the same comprehensive search strategy"
    }

@app.post("/api/presentations/generate")
async def generate_presentation_endpoint(
    request: PresentationGenerateRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate presentation with AI Slide Generator (Z.AI Style)
    
    7-Step Thinking Process:
    1. Content Analysis - การรับและวิเคราะห์ข้อมูลเบื้องต้น
    2. Planning & Direction - การวางแผนและกำหนดทิศทาง
    3. Research & Data - การค้นคว้าและรวบรวมข้อมูลเพิ่มเติม
    4. Image Generation - การค้นหารูปภาพที่เกี่ยวข้อง
    5. Create Slides - การออกแบบและสร้างสไลด์
    6. Review & Improve - การตรวจสอบและปรับปรุง
    7. Finalize & Export PPTX - การสรุปและส่งมอบผลงาน
    
    Returns both JSON and PPTX file!
    """
    try:
        print(f"\n{'='*60}")
        print(f"� AI Slide Generator (Z.AI Style)")
        print(f"{'='*60}")
        print(f"📌 Topic: {request.topic}")
        print(f"📊 Slides: {request.slide_count}")
        print(f"🎨 Style: {request.style}")
        print(f"🖼️ Images: {request.generate_images}")
        print(f"📝 Research Context: {bool(request.researchContext)}")
        
        # Use provided research context or find from blogs
        research_context = request.researchContext or ""
        
        if not research_context:
            # Check if there are recent research blogs related to this topic
            user_blog_dir = get_user_folder(request.user_id, "research_blogs")
            
            if user_blog_dir.exists():
                blogs = sorted(user_blog_dir.glob("*.json"), key=lambda x: x.stat().st_mtime, reverse=True)
                for blog_path in blogs[:5]:
                    try:
                        with open(blog_path, 'r', encoding='utf-8') as f:
                            blog_data = json.load(f)
                            if any(word.lower() in request.topic.lower() for word in blog_data['query'].split()):
                                research_context += f"\n\n### Research Context from: {blog_data['query']}\n"
                                research_context += f"{blog_data.get('executive_summary', '')}\n"
                                research_context += f"**Key Findings:**\n"
                                for finding in blog_data.get('key_findings', [])[:5]:
                                    research_context += f"- {finding}\n"
                                break
                    except Exception as e:
                        print(f"Error reading blog: {e}")
                        continue
        
        # 🚀 PRIMARY: Use AI Slide Generator (Z.AI Style with PPTX)
        result = None
        pptx_path = None
        model_used = "GLM-4.7 (Z.AI 7-Step)"
        
        if AI_SLIDE_AVAILABLE:
            try:
                print("🚀 Using AI Slide Generator (Z.AI Style + PPTX)...")
                
                result = generate_ai_slides(
                    topic=request.topic,
                    slide_count=request.slide_count,
                    style=request.style,
                    research_context=research_context if research_context else None,
                    generate_images=request.generate_images,
                    user_id=request.user_id
                )
                
                if result and result.get("success"):
                    pptx_path = result.get("pptx_path")
                    print(f"✅ Generated {len(result.get('slides', []))} slides with PPTX!")
                else:
                    result = None
                    
            except Exception as e:
                print(f"AI Slide Generator failed: {e}")
                traceback.print_exc()
                result = None
        
        # FALLBACK: Use GLM 7-Step Generator (no PPTX)
        if result is None and GLM_7STEP_AVAILABLE:
            try:
                print("📝 Fallback to GLM-4.7 7-Step Generator...")
                
                generator = GLM7StepGenerator()
                fast_result = generator.generate_presentation_fast(
                    topic=request.topic,
                    slide_count=request.slide_count,
                    style=request.style,
                    research_context=research_context if research_context else None,
                    generate_images=request.generate_images
                )
                
                if fast_result and fast_result.get("success"):
                    result = fast_result
                    model_used = "GLM-4.7 (Fallback)"
                    print(f"✅ Generated {len(fast_result.get('slides', []))} slides")
                
            except Exception as e:
                print(f"GLM-4.7 fallback failed: {e}")
                traceback.print_exc()
        
        # No result - error
        if result is None or not result.get('success'):
            raise HTTPException(
                status_code=503,
                detail="Slide generation failed. Please check GLM_API_KEY configuration."
            )
        
        # Get user presentations directory
        user_pres_dir = get_user_folder(request.user_id, "presentations")
        
        # Prepare presentation data for JSON
        presentation_data = {
            "title": result.get("title", request.topic),
            "subtitle": result.get("subtitle", ""),
            "slides": result.get("slides", []),
            "pptx_path": pptx_path,
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "topic": request.topic,
                "model": model_used,
                "style": request.style,
                "slide_count": len(result.get("slides", [])),
                "used_research_context": bool(research_context),
                "process": "7-step",
                "has_pptx": bool(pptx_path),
                "usage": result.get("metadata", {}).get("usage", {}) if isinstance(result.get("metadata"), dict) else {}
            }
        }
        
        # Save JSON to user's presentations directory
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_topic = "".join(c for c in request.topic[:50] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
        presentation_filename = f"{timestamp}_{safe_topic}.json"
        presentation_path = user_pres_dir / presentation_filename
        
        with open(presentation_path, 'w', encoding='utf-8') as f:
            json.dump(presentation_data, f, indent=2, ensure_ascii=False)
        
        print(f"✅ Presentation JSON saved: {presentation_path}")
        if pptx_path:
            print(f"✅ Presentation PPTX saved: {pptx_path}")
        
        # Generate PPTX download URL
        pptx_url = None
        if pptx_path:
            pptx_filename = Path(pptx_path).name
            pptx_url = f"/api/presentations/download/{pptx_filename}"
        
        return {
            "success": True,
            "presentation_id": presentation_filename,
            "title": presentation_data["title"],
            "subtitle": presentation_data.get("subtitle", ""),
            "slides": presentation_data["slides"],
            "pptx_path": pptx_path,
            "pptx_url": pptx_url,
            "metadata": presentation_data["metadata"]
        }
    
    except Exception as e:
        print(f"❌ Generation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

# PPTX Download Endpoint
@app.get("/api/presentations/download/{filename}")
async def download_presentation_pptx(filename: str):
    """Download generated PPTX file"""
    try:
        # Check in exports directory
        exports_dir = Path(__file__).parent / "exports" / "presentations"
        filepath = exports_dir / filename
        
        if filepath.exists() and filepath.suffix == ".pptx":
            return FileResponse(
                path=str(filepath),
                filename=filename,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        
        # Check in user_data directories
        user_data_dir = Path(__file__).parent / "user_data"
        for user_dir in user_data_dir.iterdir():
            if user_dir.is_dir():
                pres_dir = user_dir / "presentations"
                if pres_dir.exists():
                    filepath = pres_dir / filename
                    if filepath.exists() and filepath.suffix == ".pptx":
                        return FileResponse(
                            path=str(filepath),
                            filename=filename,
                            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
        
        raise HTTPException(status_code=404, detail="PPTX file not found")
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/research-blogs")
async def list_research_blogs_alt(
    user_id: str = "anonymous",
    api_key: str = Depends(verify_api_key_optional),
    limit: int = 20
):
    """List available research blogs for user"""
    try:
        # Try user-specific folder first
        user_blog_dir = get_user_folder(user_id, "research_blogs")
        
        # Also check global folder
        global_blog_dir = Path("research_blogs")
        
        blogs = []
        
        # Get blogs from user folder
        if user_blog_dir.exists():
            blog_files = sorted(user_blog_dir.glob("*.json"), key=lambda x: x.stat().st_mtime, reverse=True)
            for blog_path in blog_files[:limit]:
                try:
                    with open(blog_path, 'r', encoding='utf-8') as f:
                        blog_data = json.load(f)
                        blogs.append({
                            "id": blog_path.name,
                            "query": blog_data.get('query', ''),
                            "timestamp": blog_data.get('timestamp', ''),
                            "model": blog_data.get('model', ''),
                            "summary": (blog_data.get('executive_summary', '') or '')[:200] + "...",
                            "chat_id": blog_data.get('chat_id', ''),
                            "source": "user",
                            "key_findings": blog_data.get('key_findings', [])[:3],
                            "total_sources": len(blog_data.get('sources', []))
                        })
                except Exception as e:
                    print(f"Error reading blog {blog_path}: {e}")
                    continue
        
        # Also get from global folder
        if global_blog_dir.exists():
            blog_files = sorted(global_blog_dir.glob("*.json"), key=lambda x: x.stat().st_mtime, reverse=True)
            for blog_path in blog_files[:limit]:
                try:
                    with open(blog_path, 'r', encoding='utf-8') as f:
                        blog_data = json.load(f)
                        blogs.append({
                            "id": blog_path.name,
                            "query": blog_data.get('query', ''),
                            "timestamp": blog_data.get('timestamp', ''),
                            "model": blog_data.get('model', ''),
                            "summary": (blog_data.get('executive_summary', '') or '')[:200] + "...",
                            "chat_id": blog_data.get('chat_id', ''),
                            "source": "global",
                            "key_findings": blog_data.get('key_findings', [])[:3],
                            "total_sources": len(blog_data.get('sources', []))
                        })
                except Exception as e:
                    print(f"Error reading blog {blog_path}: {e}")
                    continue
        
        # Sort by timestamp and deduplicate
        blogs = sorted(blogs, key=lambda x: x.get('timestamp', ''), reverse=True)[:limit]
        
        print(f"📚 Found {len(blogs)} research blogs for user {user_id}")
        return {"blogs": blogs, "total": len(blogs)}
    except Exception as e:
        print(f"Error listing blogs: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/research-blogs/{blog_id}")
async def get_research_blog_alt(
    blog_id: str,
    user_id: str = "anonymous",
    api_key: str = Depends(verify_api_key_optional)
):
    """Get full research blog content"""
    try:
        # Try user folder first
        user_blog_dir = get_user_folder(user_id, "research_blogs")
        blog_path = user_blog_dir / blog_id
        
        # Fallback to global folder
        if not blog_path.exists():
            blog_path = Path("research_blogs") / blog_id
        
        if not blog_path.exists():
            raise HTTPException(status_code=404, detail="Blog not found")
        
        with open(blog_path, 'r', encoding='utf-8') as f:
            blog_data = json.load(f)
        
        return {
            "success": True,
            "blog": blog_data
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error getting blog: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ========================================
# GENERATE SLIDES FROM RESEARCH BLOG
# ========================================

class BlogToSlidesRequest(BaseModel):
    blog_id: str
    user_id: str = "anonymous"
    slide_count: int = 8
    style: str = "professional"


@app.post("/api/research-blogs/{blog_id}/generate-slides")
async def generate_slides_from_blog(
    blog_id: str,
    request: BlogToSlidesRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """Generate presentation slides from research blog content using GLM to summarize chat context first"""
    try:
        print(f"📊 Generating slides from blog: {blog_id}")
        print(f"👤 User: {request.user_id}")
        
        # Load blog content
        user_blog_dir = get_user_folder(request.user_id, "research_blogs")
        blog_path = user_blog_dir / blog_id
        
        if not blog_path.exists():
            blog_path = Path("research_blogs") / blog_id
        
        if not blog_path.exists():
            raise HTTPException(status_code=404, detail="Blog not found")
        
        with open(blog_path, 'r', encoding='utf-8') as f:
            blog_data = json.load(f)
        
        # Extract research content
        query = blog_data.get('query', '')
        executive_summary = blog_data.get('executive_summary', '')
        key_findings = blog_data.get('key_findings', [])
        recommendations = blog_data.get('recommendations', [])
        sources = blog_data.get('sources', [])
        markdown_report = blog_data.get('markdown_report', '')
        chat_id = blog_data.get('chat_id', '')
        
        # ═══════════════════════════════════════════════════════════════
        # 🔥 NEW: สรุปเนื้อหาจาก Chat Session ด้วย GLM ก่อน
        # ═══════════════════════════════════════════════════════════════
        chat_summary = ""
        important_links = []
        key_terms = []
        
        if chat_id:
            print(f"💬 Loading chat session: {chat_id}")
            try:
                # Load chat history from database or file
                from database.memory_db import MemoryDB
                db = MemoryDB()
                
                # Get chat messages
                messages = db.get_messages(chat_id)
                
                if messages:
                    # Build chat context
                    chat_context = ""
                    for msg in messages[-20:]:  # Last 20 messages
                        role = msg.get('role', 'user')
                        content = msg.get('content', '')
                        chat_context += f"\n{role.upper()}: {content[:500]}\n"
                    
                    # Use GLM to summarize chat + extract important info
                    print(f"🤖 Using GLM to summarize chat context...")
                    from GLM_core import GLMCore
                    
                    glm = GLMCore()
                    summary_prompt = f"""Analyze this chat conversation and research data to create a comprehensive summary for a presentation.

CHAT CONVERSATION:
{chat_context}

RESEARCH DATA:
Topic: {query}
Summary: {executive_summary}
Findings: {', '.join(key_findings[:5])}

Your task:
1. **Summarize the main discussion points** from the chat
2. **Extract ALL important links/URLs** mentioned (preserve exact URLs)
3. **Identify key terms and technical concepts** that MUST be included in slides
4. **Highlight critical statistics or data** from research

Return JSON format:
{{
  "chat_summary": "3-4 sentence summary of the conversation",
  "important_links": ["url1", "url2", ...],
  "key_terms": ["term1", "term2", ...],
  "critical_points": ["point1", "point2", ...],
  "data_highlights": ["stat1", "stat2", ...]
}}"""
                    
                    response = glm.chat(
                        messages=[{"role": "user", "content": summary_prompt}],
                        temperature=0.3,
                        max_tokens=2000
                    )
                    
                    # Parse GLM response
                    import re
                    content = response.get('content', '')
                    
                    # Extract JSON from response
                    json_match = re.search(r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}', content, re.DOTALL)
                    if json_match:
                        summary_data = json.loads(json_match.group(0))
                        chat_summary = summary_data.get('chat_summary', '')
                        important_links = summary_data.get('important_links', [])
                        key_terms = summary_data.get('key_terms', [])
                        critical_points = summary_data.get('critical_points', [])
                        data_highlights = summary_data.get('data_highlights', [])
                        
                        print(f"✅ GLM Summary Complete:")
                        print(f"   📝 Chat Summary: {len(chat_summary)} chars")
                        print(f"   🔗 Links Extracted: {len(important_links)}")
                        print(f"   🔑 Key Terms: {len(key_terms)}")
                    else:
                        chat_summary = content[:500]
                        
            except Exception as chat_error:
                print(f"⚠️ Chat summary error (continuing anyway): {chat_error}")
        
        # ═══════════════════════════════════════════════════════════════
        # Build ENHANCED research context with chat insights
        # ═══════════════════════════════════════════════════════════════
        research_context = f"""
# Research Topic: {query}

## Chat Discussion Context
{chat_summary}

## Key Terms to Include
{', '.join(key_terms) if key_terms else 'N/A'}

## Important Links/References
{chr(10).join(f"- {link}" for link in important_links[:10]) if important_links else 'No links extracted'}

## Executive Summary
{executive_summary}

## Key Findings
{chr(10).join(f"- {finding}" for finding in key_findings)}

## Recommendations
{chr(10).join(f"- {rec}" for rec in recommendations)}

## Data Sources ({len(sources)} references)
{chr(10).join(f"- [{s.get('title', 'Source')}]({s.get('url', '#')}) - {s.get('snippet', '')[:100]}" for s in sources[:15])}

## Detailed Analysis
{markdown_report[:4000]}
"""
        
        # Use GLM7StepGenerator for slide generation
        slide_result = None
        generator_used = None
        
        print(f"🎯 Using GLM7StepGenerator for slide generation")
        
        # Use GLM7StepGenerator (Primary)
        try:
            if GLM_7STEP_AVAILABLE:
                generator = GLM7StepGenerator()
                slide_result = generator.generate_presentation_fast(
                    topic=query,
                    slide_count=request.slide_count,
                    style=request.style,
                    research_context=research_context[:5000] if research_context else None
                )
                generator_used = "GLM-4.7 7-Step"
                print(f"✅ Slide generation using GLM7StepGenerator")
                
                # 🔥 Inject important links and key terms into slides
                if slide_result and slide_result.get('success'):
                    slides = slide_result.get('slides', [])
                    
                    # Add links to reference slide if not exists
                    has_reference_slide = any(s.get('type') == 'references' or 'reference' in s.get('title', '').lower() for s in slides)
                    
                    if important_links and not has_reference_slide:
                        slides.append({
                            "slide_number": len(slides) + 1,
                            "type": "references",
                            "title": "Important Links & References",
                            "content": important_links[:8],
                            "notes": "Key resources mentioned in the discussion"
                        })
                    
                    # Ensure key terms are visible in slides
                    if key_terms and len(slides) > 1:
                        # Add key terms to second slide if it's content
                        second_slide = slides[1] if len(slides) > 1 else None
                        if second_slide and second_slide.get('type') == 'content':
                            content = second_slide.get('content', [])
                            if isinstance(content, list):
                                content.append(f"🔑 Key Terms: {', '.join(key_terms[:5])}")
                                second_slide['content'] = content
                    
                    slide_result['slides'] = slides
                    
        except Exception as glm_error:
            print(f"⚠️ GLM slide generation failed: {glm_error}")
            import traceback
            print(traceback.format_exc())
        
        # Return if any generator succeeded
        if slide_result and slide_result.get('success'):
            return {
                "success": True,
                "title": slide_result.get('title', query),
                "slides": slide_result.get('slides', []),
                "blog_id": blog_id,
                "metadata": {
                    "generated_from": "research_blog_with_chat_context",
                    "generator": generator_used,
                    "blog_query": query,
                    "slide_count": len(slide_result.get('slides', [])),
                    "style": request.style,
                    "sources_used": len(sources),
                    "chat_summary_included": bool(chat_summary),
                    "links_extracted": len(important_links),
                    "key_terms_count": len(key_terms)
                }
            }
        
        # Fallback: Generate basic slides from content
        slides = []
        
        # Title slide
        slides.append({
            "id": "1",
            "type": "title",
            "title": query,
            "subtitle": f"Research Analysis • {len(sources)} Sources"
        })
        
        # Executive Summary slide
        slides.append({
            "id": "2",
            "type": "content",
            "title": "Executive Summary",
            "content": [executive_summary[:200] + "..." if len(executive_summary) > 200 else executive_summary]
        })
        
        # Key Findings slides
        for i, finding in enumerate(key_findings[:5], 3):
            slides.append({
                "id": str(i),
                "type": "content",
                "title": f"Key Finding {i-2}",
                "content": [finding]
            })
        
        # Recommendations slide
        if recommendations:
            slides.append({
                "id": str(len(slides) + 1),
                "type": "content",
                "title": "Recommendations",
                "content": recommendations[:4]
            })
        
        # Conclusion
        slides.append({
            "id": str(len(slides) + 1),
            "type": "conclusion",
            "title": "Conclusion",
            "content": [
                f"Analysis based on {len(sources)} credible sources",
                f"{len(key_findings)} key findings identified",
                "Thank you"
            ]
        })
        
        return {
            "success": True,
            "title": query,
            "slides": slides,
            "blog_id": blog_id,
            "metadata": {
                "generated_from": "research_blog_fallback",
                "blog_query": query,
                "slide_count": len(slides),
                "style": request.style,
                "sources_used": len(sources)
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error generating slides from blog: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))


# ========================================
# PRESENTATION MANAGEMENT ENDPOINTS
# ========================================

@app.post("/api/presentations/save")
async def save_presentation(
    request: PresentationSaveRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """Save a presentation manually"""
    try:
        print(f"💾 Saving presentation for user: {request.user_id}")
        
        # ✅ Use user-specific folder
        user_pres_dir = get_user_folder(request.user_id, "presentations")
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c for c in request.title[:50] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
        presentation_filename = f"{timestamp}_{safe_title}.json"
        presentation_path = user_pres_dir / presentation_filename
        
        presentation_data = {
            "title": request.title,
            "slides": request.slides,
            "metadata": {
                **request.metadata,
                "saved_at": datetime.now().isoformat(),
                "filename": presentation_filename
            }
        }
        
        with open(presentation_path, 'w', encoding='utf-8') as f:
            json.dump(presentation_data, f, indent=2, ensure_ascii=False)
        
        print(f"💾 Manually saved presentation: {presentation_path}")
        
        return {
            "success": True,
            "presentation_id": presentation_filename,
            "message": "Presentation saved successfully"
        }
    except Exception as e:
        print(f"❌ Save presentation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/presentations")
async def list_presentations(
    user_id: Optional[str] = "anonymous",
    api_key: str = Depends(verify_api_key_optional),
    limit: int = 20
):
    """List all saved presentations for specific user"""
    try:
        # ✅ Use user-specific folder
        user_pres_dir = get_user_folder(user_id or "anonymous", "presentations")
        
        if not user_pres_dir.exists():
            return {"presentations": []}
        
        presentations = []
        presentation_files = sorted(
            user_pres_dir.glob("*.json"), 
            key=lambda x: x.stat().st_mtime, 
            reverse=True
        )
        
        for pres_path in presentation_files[:limit]:
            try:
                with open(pres_path, 'r', encoding='utf-8') as f:
                    pres_data = json.load(f)
                    presentations.append({
                        "id": pres_path.name,
                        "title": pres_data.get('title', 'Untitled'),
                        "slide_count": len(pres_data.get('slides', [])),
                        "created_at": pres_data.get('metadata', {}).get('generated_at', ''),
                        "style": pres_data.get('metadata', {}).get('style', 'unknown'),
                        "topic": pres_data.get('metadata', {}).get('topic', '')
                    })
            except Exception as e:
                print(f"Error reading presentation {pres_path}: {e}")
                continue
        
        return {
            "success": True,
            "presentations": presentations,
            "total": len(presentations)
        }
    except Exception as e:
        print(f"❌ List presentations error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/presentations/{presentation_id}")
async def get_presentation(
    presentation_id: str,
    user_id: str,
    api_key: str = Depends(verify_api_key_optional)
):
    """Get full presentation data for specific user"""
    try:
        # ✅ Use user-specific folder
        user_pres_dir = get_user_folder(user_id, "presentations")
        presentation_path = user_pres_dir / presentation_id
        
        if not presentation_path.exists():
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        with open(presentation_path, 'r', encoding='utf-8') as f:
            presentation_data = json.load(f)
        
        return {
            "success": True,
            "presentation": presentation_data
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Get presentation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/presentations/{presentation_id}")
async def delete_presentation(
    presentation_id: str,
    api_key: str = Depends(verify_api_key_optional)
):
    """Delete a presentation"""
    try:
        presentation_path = PRESENTATIONS_DIR / presentation_id
        if not presentation_path.exists():
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        presentation_path.unlink()
        print(f"🗑️ Deleted presentation: {presentation_path}")
        
        return {
            "success": True,
            "message": "Presentation deleted successfully"
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Delete presentation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/presentations/export")
async def export_presentation_endpoint(
    request: PresentationExportRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Export presentation to PPTX or PDF
    """
    if not PRESENTATION_ENGINE_AVAILABLE:
        print("❌ Presentation engine not available")
        raise HTTPException(
            status_code=503,
            detail="Presentation engine not available"
        )
    
    if request.format not in ["pptx", "pdf"]:
        print(f"❌ Invalid format: {request.format}")
        raise HTTPException(
            status_code=400,
            detail="Format must be 'pptx' or 'pdf'"
        )
    
    try:
        print(f"📄 Exporting presentation to {request.format.upper()}")
        print(f"   Title: {request.presentation.get('title', 'Untitled')}")
        print(f"   Slides: {len(request.presentation.get('slides', []))}")
        
        # Validate presentation data
        if not request.presentation:
            raise ValueError("Presentation data is empty")
        
        if 'slides' not in request.presentation or not request.presentation['slides']:
            raise ValueError("No slides found in presentation")
        
        file_path = export_presentation(
            presentation_data=request.presentation,
            format=request.format
        )
        
        print(f"✅ Export completed: {file_path}")
        
        # Determine media type
        if request.format == "pptx":
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            media_type = "application/pdf"
        
        return FileResponse(
            path=file_path,
            media_type=media_type,
            filename=file_path.name
        )
    
    except ValueError as e:
        print(f"❌ Validation error: {str(e)}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(f"❌ Export error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════════
# Presentation Images API
# ═══════════════════════════════════════════════════════════════════
# ADMIN ENDPOINTS
# ═══════════════════════════════════════════════════════════════════

from database.db_manager import get_db

# Admin authorization helper
def verify_admin(user_email: str = Header(None, alias="X-User-Email")):
    """Verify user is admin"""
    if not user_email:
        raise HTTPException(status_code=401, detail="Missing user email header")
    
    db = get_db()
    user = db.get_user_by_email(user_email)
    
    if not user or user['role'] != 'admin':
        raise HTTPException(status_code=403, detail="Admin access required")
    
    return user

@app.get("/api/admin/stats")
async def get_admin_stats(admin_user: Dict = Depends(verify_admin)):
    """Get system statistics (admin only)"""
    try:
        db = get_db()
        stats = db.get_system_stats()
        
        # Log audit
        db.log_audit(
            user_id=admin_user['id'],
            action='view_stats',
            entity_type='system'
        )
        
        return {"success": True, "stats": stats}
    except Exception as e:
        print(f"❌ Error getting stats: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/admin/users")
async def list_users(
    role: Optional[str] = None,
    admin_user: Dict = Depends(verify_admin)
):
    """List all users (admin only)"""
    try:
        db = get_db()
        users = db.list_all_users(role=role)
        
        # Enrich with subscription info
        for user in users:
            subscription = db.get_user_subscription(user['id'])
            user['subscription'] = subscription
        
        # Log audit
        db.log_audit(
            user_id=admin_user['id'],
            action='list_users',
            entity_type='users'
        )
        
        return {"success": True, "users": users}
    except Exception as e:
        print(f"❌ Error listing users: {e}")
        raise HTTPException(status_code=500, detail=str(e))

class UserRoleUpdateRequest(BaseModel):
    user_id: str
    role: str

@app.post("/api/admin/users/role")
async def update_user_role(
    request: UserRoleUpdateRequest,
    admin_user: Dict = Depends(verify_admin)
):
    """Update user role (admin only)"""
    try:
        db = get_db()
        success = db.update_user_role(request.user_id, request.role)
        
        if success:
            # Log audit
            db.log_audit(
                user_id=admin_user['id'],
                action='update_user_role',
                entity_type='user',
                entity_id=request.user_id,
                details={'role': request.role}
            )
            
            return {"success": True, "message": "Role updated"}
        else:
            raise HTTPException(status_code=404, detail="User not found")
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Error updating role: {e}")
        raise HTTPException(status_code=500, detail=str(e))

class SubscriptionAssignRequest(BaseModel):
    user_id: str
    plan_name: str
    billing_cycle: str = "monthly"
    expires_at: Optional[str] = None

@app.post("/api/admin/subscriptions/assign")
async def assign_subscription(
    request: SubscriptionAssignRequest,
    admin_user: Dict = Depends(verify_admin)
):
    """Assign subscription to user (admin only)"""
    try:
        from datetime import datetime
        
        db = get_db()
        
        # Parse expiry date if provided
        expires_at = None
        if request.expires_at:
            expires_at = datetime.fromisoformat(request.expires_at)
        
        subscription = db.assign_subscription(
            user_id=request.user_id,
            plan_name=request.plan_name,
            billing_cycle=request.billing_cycle,
            expires_at=expires_at
        )
        
        # Log audit
        db.log_audit(
            user_id=admin_user['id'],
            action='assign_subscription',
            entity_type='subscription',
            entity_id=subscription['id'],
            details={
                'target_user': request.user_id,
                'plan': request.plan_name,
                'billing_cycle': request.billing_cycle
            }
        )
        
        return {"success": True, "subscription": subscription}
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(f"❌ Error assigning subscription: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/admin/plans")
async def get_all_plans(admin_user: Dict = Depends(verify_admin)):
    """Get all subscription plans (admin only)"""
    try:
        db = get_db()
        plans = db.get_all_plans()
        return {"success": True, "plans": plans}
    except Exception as e:
        print(f"❌ Error getting plans: {e}")
        raise HTTPException(status_code=500, detail=str(e))

class WhitelistAddRequest(BaseModel):
    email: str
    plan_name: Optional[str] = None
    notes: Optional[str] = None

@app.post("/api/admin/whitelist/add")
async def add_to_whitelist(
    request: WhitelistAddRequest,
    admin_user: Dict = Depends(verify_admin)
):
    """Add email to whitelist (admin only)"""
    try:
        db = get_db()
        entry = db.add_to_whitelist(
            email=request.email,
            plan_name=request.plan_name,
            added_by=admin_user['id'],
            notes=request.notes
        )
        
        # Log audit
        db.log_audit(
            user_id=admin_user['id'],
            action='add_whitelist',
            entity_type='whitelist',
            entity_id=entry['id'],
            details={'email': request.email, 'plan': request.plan_name}
        )
        
        return {"success": True, "entry": entry}
    except Exception as e:
        print(f"❌ Error adding to whitelist: {e}")
        raise HTTPException(status_code=500, detail=str(e))

class WhitelistRemoveRequest(BaseModel):
    email: str

@app.post("/api/admin/whitelist/remove")
async def remove_from_whitelist(
    request: WhitelistRemoveRequest,
    admin_user: Dict = Depends(verify_admin)
):
    """Remove email from whitelist (admin only)"""
    try:
        db = get_db()
        success = db.remove_from_whitelist(request.email)
        
        if success:
            # Log audit
            db.log_audit(
                user_id=admin_user['id'],
                action='remove_whitelist',
                entity_type='whitelist',
                details={'email': request.email}
            )
            
            return {"success": True, "message": "Removed from whitelist"}
        else:
            raise HTTPException(status_code=404, detail="Email not found in whitelist")
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Error removing from whitelist: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/admin/whitelist")
async def list_whitelist(
    is_active: Optional[bool] = None,
    admin_user: Dict = Depends(verify_admin)
):
    """List whitelisted emails (admin only)"""
    try:
        db = get_db()
        entries = db.list_whitelist(is_active=is_active)
        return {"success": True, "entries": entries}
    except Exception as e:
        print(f"❌ Error listing whitelist: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/admin/audit-logs")
async def get_audit_logs(
    user_id: Optional[str] = None,
    limit: int = 100,
    admin_user: Dict = Depends(verify_admin)
):
    """Get audit logs (admin only)"""
    try:
        db = get_db()
        logs = db.get_audit_logs(user_id=user_id, limit=limit)
        return {"success": True, "logs": logs}
    except Exception as e:
        print(f"❌ Error getting audit logs: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/admin/usage/{user_id}")
async def get_user_usage_admin(
    user_id: str,
    admin_user: Dict = Depends(verify_admin)
):
    """Get user usage statistics (admin only)"""
    try:
        from datetime import datetime
        
        db = get_db()
        
        # Get current month usage
        now = datetime.now()
        monthly_usage = db.get_monthly_usage(user_id, now.year, now.month)
        
        # Get subscription and limits
        subscription = db.get_user_subscription(user_id)
        
        return {
            "success": True,
            "usage": monthly_usage,
            "subscription": subscription,
            "period": {"year": now.year, "month": now.month}
        }
    except Exception as e:
        print(f"❌ Error getting user usage: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════════
# USER ENDPOINTS (for self-service)
# ═══════════════════════════════════════════════════════════════════

@app.get("/api/user/subscription")
async def get_my_subscription(user_email: str = Header(None, alias="X-User-Email")):
    """Get current user's subscription"""
    try:
        if not user_email:
            raise HTTPException(status_code=401, detail="Missing user email")
        
        db = get_db()
        user = db.get_or_create_user(user_email)
        subscription = db.get_user_subscription(user['id'])
        
        # If no subscription, assign default plan from whitelist
        if not subscription:
            whitelist = db.check_whitelist(user_email)
            if whitelist and whitelist.get('plan_name'):
                subscription = db.assign_subscription(
                    user_id=user['id'],
                    plan_name=whitelist['plan_name'],
                    billing_cycle='monthly'
                )
            else:
                # Assign free Plus plan by default
                subscription = db.assign_subscription(
                    user_id=user['id'],
                    plan_name='Plus',
                    billing_cycle='monthly'
                )
        
        return {"success": True, "subscription": subscription, "user": user}
    except Exception as e:
        print(f"❌ Error getting subscription: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/user/usage")
async def get_my_usage(user_email: str = Header(None, alias="X-User-Email")):
    """Get current user's usage statistics"""
    try:
        if not user_email:
            raise HTTPException(status_code=401, detail="Missing user email")
        
        from datetime import datetime
        
        db = get_db()
        user = db.get_or_create_user(user_email)
        
        # Get current month usage
        now = datetime.now()
        monthly_usage = db.get_monthly_usage(user['id'], now.year, now.month)
        
        # Get subscription and limits
        subscription = db.get_user_subscription(user['id'])
        
        # Check limits for each resource type
        limits_check = {}
        resource_types = ['message', 'token', 'research', 'presentation', 'image']
        for resource_type in resource_types:
            limits_check[resource_type] = db.check_usage_limit(user['id'], resource_type)
        
        return {
            "success": True,
            "usage": monthly_usage,
            "limits": limits_check,
            "subscription": subscription,
            "period": {"year": now.year, "month": now.month}
        }
    except Exception as e:
        print(f"❌ Error getting usage: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/plans")
async def get_public_plans():
    """Get all available subscription plans (public)"""
    try:
        db = get_db()
        plans = db.get_all_plans()
        return {"success": True, "plans": plans}
    except Exception as e:
        print(f"❌ Error getting plans: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════════

@app.head("/api/presentations/images/{filename}")
@app.get("/api/presentations/images/{filename}")
async def get_presentation_image(
    filename: str,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Serve generated presentation images
    GET /api/presentations/images/{filename}
    HEAD /api/presentations/images/{filename}
    """
    try:
        # Check if file exists
        images_dir = Path("exports/presentations/images")
        image_path = images_dir / filename
        
        if not image_path.exists():
            raise HTTPException(status_code=404, detail="Image not found")
        
        # Verify it's an image file
        if not filename.lower().endswith(('.png', '.jpg', '.jpeg', '.webp')):
            raise HTTPException(status_code=400, detail="Invalid file type")
        
        # Return image file
        return FileResponse(
            image_path,
            media_type="image/png" if filename.endswith('.png') else "image/jpeg",
            headers={
                "Cache-Control": "public, max-age=3600",
                "Access-Control-Allow-Origin": "*"
            }
        )
    
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Image serve error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════════
# Z-STYLE SLIDE GENERATION (2-Step with Streaming)
# ═══════════════════════════════════════════════════════════════════

class OutlineRequest(BaseModel):
    """Request for Step 1: Generate Outline"""
    topic: str
    slide_count: int = 8
    style: str = "professional"
    model: str = "auto"  # "glm", "azure", or "auto"
    user_id: Optional[str] = "anonymous"

class FullSlidesRequest(BaseModel):
    """Request for Step 2: Generate Full Slides"""
    outline: Dict[str, Any]
    generate_images: bool = True
    model: str = "auto"  # "glm", "azure", or "auto"
    research_context: Optional[str] = None
    user_id: Optional[str] = "anonymous"

@app.post("/api/zstyle/outline")
async def zstyle_generate_outline(
    request: OutlineRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Z-Style 7-Step Generation: Steps 1-2 (Content Analysis & Planning)
    Returns presentation outline with design plan
    
    Uses GLM-4.7 via Z.AI Direct API (cheaper!)
    """
    try:
        from engines.glm_7step_generator import get_generator
        
        print(f"📋 [7-Step] Generating outline: {request.topic}")
        print(f"   🤖 Model: GLM-4.7 (Z.AI Direct API)")
        
        generator = get_generator()
        
        # Run Steps 1-2 and return outline
        result = {
            "success": True,
            "step": "outline",
            "topic": request.topic,
            "model_used": "glm-4.7",
            "data": None
        }
        
        # Step 1: Content Analysis
        try:
            analysis = generator._step1_content_analysis(request.topic, None)
            result["analysis"] = analysis
        except Exception as e:
            print(f"⚠️ Analysis error: {e}")
            analysis = {"content_type": "general"}
        
        # Step 2: Planning
        try:
            plan = generator._step2_planning(
                request.topic, request.slide_count, request.style, analysis
            )
            result["data"] = {
                "title": request.topic,
                "subtitle": plan.get("content_structure", {}).get("opening_type", "") if isinstance(plan.get("content_structure"), dict) else "",
                "outline": plan.get("slide_outline", []),
                "design_plan": plan.get("design_plan", {}),
                "needs_research": plan.get("needs_research", False)
            }
            result["message"] = "Outline generated with 7-step process! Review and confirm."
        except Exception as e:
            print(f"⚠️ Planning error: {e}")
            # Fallback to default outline
            result["data"] = {
                "title": request.topic,
                "outline": generator._generate_default_outline(request.slide_count)
            }
            result["message"] = "Outline generated (fallback). Review and confirm."
        
        return result
        
    except Exception as e:
        print(f"❌ Outline generation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/zstyle/slides/stream")
async def zstyle_generate_slides_stream(
    request: FullSlidesRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Z-Style 7-Step Generation: Steps 3-7 with Server-Sent Events (SSE)
    Streams progress for each step and returns slides one by one
    
    Uses GLM-4.7 via Z.AI Direct API (cheaper!)
    """
    try:
        from engines.glm_7step_generator import get_generator
        import asyncio
        
        print(f"🎨 [7-Step] Streaming slide generation...")
        print(f"   🤖 Model: GLM-4.7 (Z.AI Direct API)")
        
        generator = get_generator()
        
        # Get outline info
        outline = request.outline
        topic = outline.get("title", "Presentation")
        slide_outline = outline.get("outline", [])
        slide_count = len(slide_outline) if slide_outline else 8
        
        async def event_generator():
            """Generate SSE events using GLM-4.7 7-step process"""
            try:
                import concurrent.futures
                import asyncio
                
                # Run generator in thread pool (it's synchronous)
                def run_generator():
                    return list(generator.generate_presentation(
                        topic=topic,
                        slide_count=slide_count,
                        style=request.outline.get("style", "professional"),
                        research_context=request.research_context,
                        generate_images=request.generate_images
                    ))
                
                loop = asyncio.get_event_loop()
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    events = await loop.run_in_executor(executor, run_generator)
                
                # Stream the events
                slides = []
                for event in events:
                    event_type = event.get("type", "")
                    
                    if event_type == "step_start":
                        yield f"data: {json.dumps({'type': 'step_start', 'step': event.get('step'), 'name': event.get('name', ''), 'message': event.get('message', '')}, ensure_ascii=False)}\n\n"
                    
                    elif event_type == "step_complete":
                        yield f"data: {json.dumps({'type': 'step_complete', 'step': event.get('step'), 'message': event.get('message', '')}, ensure_ascii=False)}\n\n"
                    
                    elif event_type == "slide":
                        slide = event.get("slide", {})
                        slides.append(slide)
                        slide_event = {
                            "type": "slide",
                            "slide_number": event.get("slide_number", len(slides)),
                            "slide": slide
                        }
                        yield f"data: {json.dumps(slide_event, ensure_ascii=False)}\n\n"
                    
                    elif event_type == "complete":
                        complete_event = {
                            "type": "complete",
                            "title": topic,
                            "slides": slides,
                            "message": event.get("message", f"🎉 สร้าง {len(slides)} สไลด์เสร็จสมบูรณ์!")
                        }
                        yield f"data: {json.dumps(complete_event, ensure_ascii=False)}\n\n"
                    
                    await asyncio.sleep(0.05)
                
                yield "data: {\"type\": \"done\"}\n\n"
                
            except Exception as e:
                print(f"❌ Stream error: {e}")
                error_event = {"type": "error", "message": str(e)}
                yield f"data: {json.dumps(error_event)}\n\n"
        
        return StreamingResponse(
            event_generator(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no",
                "Access-Control-Allow-Origin": "*"
            }
        )
        
    except Exception as e:
        print(f"❌ Stream generation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/zstyle/slides")
async def zstyle_generate_slides_full(
    request: FullSlidesRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Z-Style 7-Step (Non-Streaming): Generate full slides at once
    Uses GLM-4.7 via Z.AI Direct API
    """
    try:
        from engines.glm_7step_generator import get_generator
        
        print(f"🎨 [7-Step] Generating full slides (non-streaming)...")
        print(f"   🤖 Model: GLM-4.7 (Z.AI Direct API)")
        generator = get_generator()
        
        # Get topic from outline
        outline = request.outline
        topic = outline.get("title", "Presentation")
        slide_count = len(outline.get("outline", [])) or 8
        
        slides = []
        summary = None
        
        for event in generator.generate_presentation(
            topic=topic,
            slide_count=slide_count,
            style=outline.get("style", "professional"),
            research_context=request.research_context,
            generate_images=request.generate_images
        ):
            if event["type"] == "slide":
                slides.append(event["slide"])
            elif event["type"] == "complete":
                summary = event.get("summary")
                
                # Save presentation
                user_pres_dir = get_user_folder(request.user_id, "presentations")
                
                presentation_data = {
                    "title": event.get("title", topic),
                    "slides": slides,
                    "metadata": {
                        "generated_at": datetime.now().isoformat(),
                        "model": "GLM-4.7 (ZhipuAI) + CogView-3",
                        "style": "7-step",
                        "slide_count": len(slides),
                        "summary": summary
                    }
                }
                
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                safe_title = "".join(c for c in event.get("title", "Untitled")[:50] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
                filename = f"{timestamp}_{safe_title}.json"
                
                with open(user_pres_dir / filename, 'w', encoding='utf-8') as f:
                    json.dump(presentation_data, f, indent=2, ensure_ascii=False)
                
                return {
                    "success": True,
                    "presentation_id": filename,
                    "title": event.get("title"),
                    "slides": slides,
                    "metadata": presentation_data["metadata"],
                    "summary": summary
                }
        
        raise HTTPException(status_code=500, detail="Generation did not complete")
        
    except Exception as e:
        print(f"❌ Full generation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════════
# CODE-BASED SLIDE GENERATOR API (Z.AI Style - HTML/CSS Slides)
# ═══════════════════════════════════════════════════════════════════

class CodeSlideRequest(BaseModel):
    """Request for code-based slide generation"""
    topic: str
    outline: List[Dict[str, Any]]
    theme: str = "professional"
    user_id: Optional[str] = "anonymous"

@app.post("/api/codeslide/generate")
async def codeslide_generate(
    request: CodeSlideRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate slides as HTML/CSS code (Z.AI style)
    Returns slides with HTML code that can be rendered in browser
    """
    try:
        from engines.code_slide_generator import get_code_slide_generator
        
        print(f"🎨 [CodeSlide] Generating HTML slides: {request.topic}")
        generator = get_code_slide_generator()
        
        result = generator.generate_slides_sync(
            topic=request.topic,
            outline=request.outline,
            theme=request.theme
        )
        
        return result
        
    except Exception as e:
        print(f"❌ CodeSlide generation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/codeslide/stream")
async def codeslide_stream(
    request: CodeSlideRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate slides as HTML/CSS code with SSE streaming
    Each slide is sent as it's generated
    """
    try:
        from engines.code_slide_generator import get_code_slide_generator
        
        print(f"🎨 [CodeSlide] Streaming HTML slides: {request.topic}")
        generator = get_code_slide_generator()
        
        async def event_generator():
            """Generate SSE events for each slide"""
            try:
                async for event in generator.generate_slides_stream(
                    topic=request.topic,
                    outline=request.outline,
                    theme=request.theme
                ):
                    data = json.dumps(event, ensure_ascii=False)
                    yield f"data: {data}\n\n"
                    
            except Exception as e:
                error_event = json.dumps({
                    "type": "error",
                    "error": str(e)
                })
                yield f"data: {error_event}\n\n"
        
        return StreamingResponse(
            event_generator(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no"
            }
        )
        
    except Exception as e:
        print(f"❌ CodeSlide streaming error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════════
# Z-STYLE SLIDE GENERATOR (GLM-4.7 or Azure GPT-4o, Beautiful HTML/CSS)
# ═══════════════════════════════════════════════════════════════════

class ZStyleRequest(BaseModel):
    """Request for Z-Style slide generation (7-Step GLM-Style Process)"""
    topic: str
    slide_count: int = 8
    theme: str = "dark"
    model: str = "glm-7step"  # GLM-4.7 7-Step Only
    user_id: Optional[str] = "anonymous"
    generate_images: bool = True  # Generate CogView-3 images
    use_research: bool = True  # Use Deep Research for content

@app.post("/api/zslides/generate")
async def zslides_generate(
    request: ZStyleRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate beautiful Z.AI-style slides using GLM-4.7 7-Step Process
    
    Model: GLM-4.7 via ZhipuAI Direct API (7-Step Process)
    Images: CogView-3 (ZhipuAI)
    
    7 Steps:
    1. Content Analysis
    2. Planning & Direction
    3. Research & Data
    4. Image Generation (CogView-3)
    5. Create Slides
    6. Review & Improve
    7. Finalize
    """
    try:
        print(f"[GLM-4.7 Z.AI 7-Step] Generating slides: {request.topic}")
        print(f"   Research: {request.use_research}, Images: {request.generate_images}")
        
        # Use GLM-4.7 7-step generator (Z.AI API)
        generator = GLM7StepGenerator()
        
        # Collect events from generator
        slides = []
        title = request.topic
        steps_completed = []
        
        for event in generator.generate_presentation(
            topic=request.topic,
            slide_count=request.slide_count,
            style=request.theme,
            generate_images=request.generate_images,
            research_context=None
        ):
            event_type = event.get("type", "")
            if event_type == "slide":
                slides.append(event.get("slide", {}))
            elif event_type == "step_complete":
                steps_completed.append(event)
            elif event_type == "complete":
                title = event.get("title", request.topic)
                # Get slides from complete event if not collected yet
                if not slides and event.get("slides"):
                    slides = event.get("slides", [])
        
        if slides:
            return {
                "success": True,
                "presentation": {
                    "title": title,
                    "slides": slides,
                    "slide_count": len(slides)
                },
                "slides": slides,
                "steps": steps_completed,
                "model": "GLM-4.7 (Z.AI 7-Step)"
            }
        else:
            raise HTTPException(status_code=500, detail="Failed to generate slides")
        
    except Exception as e:
        print(f"GLM 7-Step generation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/zslides/generate/stream")
async def zslides_generate_stream(
    request: ZStyleRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate Z.AI-style slides with real-time SSE streaming
    
    Model: GLM-4.7 via ZhipuAI Direct API (7-Step Process)
    Images: CogView-3 (ZhipuAI)
    
    Returns SSE stream with step updates and slides as they are generated.
    """
    try:
        print(f"[GLM-4.7 Z.AI 7-Step Stream] Generating slides: {request.topic}")
        print(f"   Research: {request.use_research}, Images: {request.generate_images}")
        
        async def event_generator():
            """Generate SSE events for real-time updates"""
            import concurrent.futures
            
            generator = GLM7StepGenerator()
            
            def run_generator():
                """Run generator in thread"""
                events = []
                for event in generator.generate_presentation(
                    topic=request.topic,
                    slide_count=request.slide_count,
                    style=request.theme,
                    generate_images=request.generate_images,
                    research_context=None
                ):
                    events.append(event)
                return events
            
            # Run generator in thread pool
            loop = asyncio.get_event_loop()
            with concurrent.futures.ThreadPoolExecutor() as executor:
                events = await loop.run_in_executor(executor, run_generator)
            
            # Stream events
            slides = []
            for event in events:
                event_type = event.get("type", "")
                
                # Send step progress
                if event_type == "step_start":
                    yield f"data: {json.dumps({'type': 'step_start', 'step': event.get('step'), 'name': event.get('name'), 'message': event.get('message')}, ensure_ascii=False)}\n\n"
                
                elif event_type == "step_complete":
                    yield f"data: {json.dumps({'type': 'step_complete', 'step': event.get('step'), 'message': event.get('message')}, ensure_ascii=False)}\n\n"
                
                elif event_type == "slide":
                    slide = event.get("slide", {})
                    slides.append(slide)
                    yield f"data: {json.dumps({'type': 'slide', 'slide_number': event.get('slide_number'), 'slide': slide}, ensure_ascii=False)}\n\n"
                
                elif event_type == "complete":
                    complete_data = {
                        "type": "complete",
                        "title": event.get("title", request.topic),
                        "slides": slides,
                        "message": f"🎉 สร้าง {len(slides)} สไลด์เสร็จสมบูรณ์!"
                    }
                    yield f"data: {json.dumps(complete_data, ensure_ascii=False)}\n\n"
                
                await asyncio.sleep(0.05)  # Small delay between events
            
            yield "data: {\"type\": \"done\"}\n\n"
        
        return StreamingResponse(
            event_generator(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no",
                "Access-Control-Allow-Origin": "*"
            }
        )
        
    except Exception as e:
        print(f"GLM 7-Step stream error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/zslides/models")
async def zslides_list_models():
    """
    List available models for slide generation - GLM-4.7 Z.AI Only
    """
    return {
        "success": True,
        "available_models": {
            "glm-4.7": GLM_7STEP_AVAILABLE,
            "dalle-3": bool(os.getenv("AZURE_OPENAI_API_KEY"))
        },
        "current_model": "GLM-4.7 (Z.AI 7-Step)",
        "model_id": "glm-4.7",
        "api": "Z.AI Direct API",
        "thinking_process": "7-step",
        "features": {
            "dalle_images": bool(os.getenv("AZURE_DALL_E_ENDPOINT")),
            "7step_workflow": True
        },
        "steps": [
            "1. Content Analysis",
            "2. Planning & Direction",
            "3. Research & Data",
            "4. Image Generation",
            "5. Create Slides",
            "6. Review & Improve",
            "7. Finalize"
        ]
    }

@app.post("/api/zslides/outline")
async def zslides_outline(
    request: ZStyleRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate outline only using GLM-4.7 7-Step (Steps 1-2)
    """
    try:
        print(f"[GLM-4.7 Z.AI] Generating outline: {request.topic}")
        generator = GLM7StepGenerator()
        
        # Just run step 1-2 to get outline
        outline_slides = []
        for event in generator.generate_presentation(
            topic=request.topic,
            slide_count=request.slide_count,
            style=request.theme,
            generate_images=False,
            research_context=None
        ):
            if event.get("type") == "step_complete" and event.get("step") == 2:
                outline_slides = event.get("result", {}).get("outline", [])
                break
        
        return {
            "success": True,
            "topic": request.topic,
            "slide_count": request.slide_count,
            "model_used": "GLM-4.7 (Z.AI)",
            "thinking_process": "7-step",
            "outline": outline_slides
        }
        
    except Exception as e:
        print(f"Outline error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/zslides/download")
async def zslides_download(
    request: ZStyleRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate and download presentation as HTML file (using GLM-4.7 7-Step)
    """
    try:
        print(f"[GLM-4.7 Z.AI] Generating for download: {request.topic}")
        generator = GLM7StepGenerator()
        
        # Generate full presentation
        slides = []
        presentation = None
        
        for event in generator.generate_presentation(
            topic=request.topic,
            slide_count=request.slide_count,
            style=request.theme,
            generate_images=getattr(request, 'generate_images', True),
            research_context=None
        ):
            if event.get("type") == "slide":
                slides.append(event.get("slide", {}))
            elif event.get("type") == "complete":
                presentation = event.get("presentation", {})
        
        if presentation and slides:
            # Generate simple HTML
            html_content = f"""<!DOCTYPE html>
<html lang="th">
<head>
    <meta charset="UTF-8">
    <title>{request.topic}</title>
    <style>
        body {{ font-family: 'Segoe UI', sans-serif; margin: 0; padding: 20px; background: #f5f5f5; }}
        .slide {{ background: white; padding: 40px; margin: 20px auto; max-width: 960px; border-radius: 10px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }}
        .slide h1 {{ color: #2563EB; margin-bottom: 20px; }}
        .slide ul {{ line-height: 1.8; }}
        img {{ max-width: 100%; border-radius: 8px; }}
    </style>
</head>
<body>
"""
            for slide in slides:
                html_content += f"""<div class="slide">
    <h1>{slide.get('title', '')}</h1>
    <ul>"""
                content = slide.get('content', [])
                if isinstance(content, list):
                    for item in content:
                        html_content += f"<li>{item}</li>"
                html_content += "</ul>"
                if slide.get('imageUrl'):
                    html_content += f'<img src="{slide["imageUrl"]}" alt="{slide.get("title", "")}">'
                html_content += "</div>"
            
            html_content += "</body></html>"
            
            safe_topic = request.topic[:50].replace(" ", "_").replace("/", "-")
            filename = f"glm_presentation_{safe_topic}.html"
            
            from fastapi.responses import Response
            return Response(
                content=html_content,
                media_type="text/html",
                headers={
                    "Content-Disposition": f'attachment; filename="{filename}"',
                    "Content-Type": "text/html; charset=utf-8"
                }
            )
        else:
            raise HTTPException(status_code=500, detail="Failed to generate presentation")
        
    except Exception as e:
        print(f"Download error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════════
# CHAT-TO-PRESENTATION API (Z.AI Style)
# ═══════════════════════════════════════════════════════════════════

class ChatToPresentationRequest(BaseModel):
    """Request for chat-to-presentation generation"""
    chat_history: List[Dict[str, str]]  # List of {"role": "user/assistant", "content": "..."}
    topic_hint: Optional[str] = None  # Optional topic to focus on
    slide_count: int = 8
    style: str = "professional"
    generate_images: bool = True
    user_id: str = "anonymous"

class ResearchBlogChatRequest(BaseModel):
    """Request to save research as chat with presentation capability"""
    title: str
    query: str
    content: str
    chat_messages: List[Dict[str, str]]  # Chat history that led to this research
    sources: List[Dict[str, Any]] = []
    metadata: Dict[str, Any] = {}
    user_id: str = "anonymous"

@app.post("/api/chat-to-presentation")
async def chat_to_presentation(
    request: ChatToPresentationRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate presentation from chat history (Z.AI Style)
    
    This endpoint:
    1. Summarizes the chat conversation
    2. Extracts key points and facts
    3. Generates slides using GLM-4.7 7-Step process
    """
    try:
        print(f"\n{'='*60}")
        print(f"💬 Chat-to-Presentation Generation")
        print(f"{'='*60}")
        print(f"📝 Messages: {len(request.chat_history)}")
        print(f"📊 Slides: {request.slide_count}")
        print(f"🎨 Style: {request.style}")
        print(f"{'='*60}\n")
        
        if not GLM_7STEP_AVAILABLE:
            raise HTTPException(
                status_code=503,
                detail="GLM-4.7 7-Step Generator not available. Check GLM_API_KEY."
            )
        
        generator = GLM7StepGenerator()
        
        # Collect events from generator
        slides = []
        summary = None
        chat_summary = None
        
        for event in generator.generate_from_chat(
            chat_history=request.chat_history,
            topic_hint=request.topic_hint,
            slide_count=request.slide_count,
            style=request.style,
            generate_images=request.generate_images
        ):
            event_type = event.get("type", "")
            
            if event_type == "step_complete" and event.get("step") == 0:
                chat_summary = event.get("result", {})
            elif event_type == "slide":
                slides.append(event.get("slide", {}))
            elif event_type == "complete":
                summary = event.get("summary", {})
        
        if slides:
            # Save the presentation
            user_pres_dir = get_user_folder(request.user_id, "presentations")
            
            topic = chat_summary.get("topic", request.topic_hint or "Chat Presentation")
            
            presentation_data = {
                "title": topic,
                "slides": slides,
                "metadata": {
                    "generated_at": datetime.now().isoformat(),
                    "model": "GLM-4.7 (Z.AI Chat-to-Slides)",
                    "style": request.style,
                    "slide_count": len(slides),
                    "source": "chat",
                    "chat_summary": chat_summary,
                    "summary": summary
                }
            }
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = "".join(c for c in topic[:50] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
            filename = f"{timestamp}_chat_{safe_title}.json"
            
            with open(user_pres_dir / filename, 'w', encoding='utf-8') as f:
                json.dump(presentation_data, f, indent=2, ensure_ascii=False)
            
            print(f"✅ Chat presentation saved: {filename}")
            
            return {
                "success": True,
                "presentation_id": filename,
                "title": topic,
                "slides": slides,
                "chat_summary": chat_summary,
                "metadata": presentation_data["metadata"]
            }
        else:
            raise HTTPException(status_code=500, detail="Failed to generate slides from chat")
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Chat-to-presentation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/chat-to-presentation/stream")
async def chat_to_presentation_stream(
    request: ChatToPresentationRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate presentation from chat history with SSE streaming
    """
    try:
        print(f"💬 [Stream] Chat-to-Presentation: {len(request.chat_history)} messages")
        
        if not GLM_7STEP_AVAILABLE:
            raise HTTPException(
                status_code=503,
                detail="GLM-4.7 7-Step Generator not available"
            )
        
        generator = GLM7StepGenerator()
        
        async def event_generator():
            """Generate SSE events"""
            import concurrent.futures
            
            def run_generator():
                return list(generator.generate_from_chat(
                    chat_history=request.chat_history,
                    topic_hint=request.topic_hint,
                    slide_count=request.slide_count,
                    style=request.style,
                    generate_images=request.generate_images
                ))
            
            loop = asyncio.get_event_loop()
            with concurrent.futures.ThreadPoolExecutor() as executor:
                events = await loop.run_in_executor(executor, run_generator)
            
            slides = []
            for event in events:
                event_type = event.get("type", "")
                
                if event_type in ["step_start", "step_complete"]:
                    yield f"data: {json.dumps(event, ensure_ascii=False)}\n\n"
                
                elif event_type == "slide":
                    slides.append(event.get("slide", {}))
                    yield f"data: {json.dumps(event, ensure_ascii=False)}\n\n"
                
                elif event_type == "complete":
                    complete_event = {
                        "type": "complete",
                        "title": event.get("title", "Presentation"),
                        "slides": slides,
                        "message": event.get("message", f"🎉 สร้าง {len(slides)} สไลด์เสร็จสมบูรณ์!")
                    }
                    yield f"data: {json.dumps(complete_event, ensure_ascii=False)}\n\n"
                
                await asyncio.sleep(0.05)
            
            yield "data: {\"type\": \"done\"}\n\n"
        
        return StreamingResponse(
            event_generator(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no",
                "Access-Control-Allow-Origin": "*"
            }
        )
        
    except Exception as e:
        print(f"❌ Stream error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/research-blogs/save-with-chat")
async def save_research_blog_with_chat(
    request: ResearchBlogChatRequest,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Save research blog along with the chat history that led to it
    This enables generating presentations from the research context later
    """
    try:
        print(f"💾 Saving research blog with chat: {request.title}")
        
        user_blog_dir = get_user_folder(request.user_id, "research_blogs")
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c for c in request.title[:50] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
        filename = f"{timestamp}_{safe_title}.json"
        
        blog_data = {
            "id": timestamp,
            "title": request.title,
            "query": request.query,
            "content": request.content,
            "chat_messages": request.chat_messages,  # Store chat history
            "sources": request.sources,
            "metadata": {
                **request.metadata,
                "saved_at": datetime.now().isoformat(),
                "filename": filename,
                "has_chat_context": True,
                "chat_message_count": len(request.chat_messages)
            }
        }
        
        with open(user_blog_dir / filename, 'w', encoding='utf-8') as f:
            json.dump(blog_data, f, indent=2, ensure_ascii=False)
        
        print(f"✅ Research blog with chat saved: {filename}")
        
        return {
            "success": True,
            "blog_id": filename,
            "message": "Research blog saved with chat context",
            "can_generate_slides": True
        }
        
    except Exception as e:
        print(f"❌ Save error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/research-blogs/{blog_id}/generate-from-chat")
async def generate_slides_from_blog_chat(
    blog_id: str,
    user_id: str = "anonymous",
    slide_count: int = 8,
    style: str = "professional",
    generate_images: bool = True,
    api_key: str = Depends(verify_api_key_optional)
):
    """
    Generate presentation from research blog using its chat context
    """
    try:
        print(f"📊 Generating slides from blog chat: {blog_id}")
        
        # Load blog with chat context
        user_blog_dir = get_user_folder(user_id, "research_blogs")
        blog_path = user_blog_dir / blog_id
        
        if not blog_path.exists():
            blog_path = Path("research_blogs") / blog_id
        
        if not blog_path.exists():
            raise HTTPException(status_code=404, detail="Blog not found")
        
        with open(blog_path, 'r', encoding='utf-8') as f:
            blog_data = json.load(f)
        
        # Check if blog has chat context
        chat_messages = blog_data.get('chat_messages', [])
        
        if not GLM_7STEP_AVAILABLE:
            raise HTTPException(
                status_code=503,
                detail="GLM-4.7 7-Step Generator not available"
            )
        
        generator = GLM7StepGenerator()
        
        if chat_messages:
            # Generate from chat context
            slides = []
            for event in generator.generate_from_chat(
                chat_history=chat_messages,
                topic_hint=blog_data.get('query', blog_data.get('title', '')),
                slide_count=slide_count,
                style=style,
                generate_images=generate_images
            ):
                if event.get("type") == "slide":
                    slides.append(event.get("slide", {}))
                elif event.get("type") == "complete":
                    break
        else:
            # Fallback to research context
            slides = []
            for event in generator.generate_from_research_blog(
                blog_content=blog_data.get('content', ''),
                blog_title=blog_data.get('title', blog_data.get('query', 'Research')),
                sources=blog_data.get('sources', []),
                slide_count=slide_count,
                style=style,
                generate_images=generate_images
            ):
                if event.get("type") == "slide":
                    slides.append(event.get("slide", {}))
                elif event.get("type") == "complete":
                    break
        
        if slides:
            return {
                "success": True,
                "title": blog_data.get('title', blog_data.get('query', 'Presentation')),
                "slides": slides,
                "blog_id": blog_id,
                "metadata": {
                    "generated_from": "research_blog_chat" if chat_messages else "research_blog",
                    "generator": "GLM-4.7 (Z.AI 7-Step)",
                    "slide_count": len(slides),
                    "style": style
                }
            }
        else:
            raise HTTPException(status_code=500, detail="Failed to generate slides")
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Generation error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(
        "backend_api:app",
        host="0.0.0.0",
        port=8000,
        reload=True,
        log_level="info"
    )
