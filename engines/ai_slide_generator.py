"""
AI Slide Generator - Z.AI Style Complete Solution
================================================================
🎯 7-Step Thinking Process สำหรับสร้าง Presentation ที่สวยงาม
🎨 สร้าง PPTX อัตโนมัติพร้อม Design ที่สวยงาม

Features:
- 7-Step Thinking: วิเคราะห์ → วางแผน → ค้นคว้า → ออกแบบ → สร้าง → ตรวจสอบ → ส่งมอบ
- Beautiful PPTX: สร้าง slide สวยงามแบบมืออาชีพ
- Single API Call: เร็วและมีประสิทธิภาพ
- Image Support: รองรับ CogView-3 และ DALL-E

Model: GLM-4.7 (ZhipuAI BigModel API)
"""
import os
import sys
import json
import requests
import time
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from io import BytesIO

# Load environment
PROJECT_ROOT = Path(__file__).parent.parent
ENV_PATH = PROJECT_ROOT / ".env"
load_dotenv(dotenv_path=ENV_PATH, override=True)

# API Keys
GLM_API_KEY = os.getenv("GLM_API_KEY", "5c6ede62c88041158aba9d710f700a0e.0bncgwYVObXu0Hyh")
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_DALLE_ENDPOINT = os.getenv("AZURE_DALL_E_ENDPOINT") or os.getenv("AZURE_DALLE_ENDPOINT")

# PowerPoint libraries
try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
    print("✅ python-pptx loaded successfully")
except ImportError as e:
    PPTX_AVAILABLE = False
    print(f"❌ python-pptx not available: {e}")

# PIL for images
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

# Exports directory
EXPORTS_DIR = PROJECT_ROOT / "exports" / "presentations"
IMAGES_DIR = EXPORTS_DIR / "images"
EXPORTS_DIR.mkdir(parents=True, exist_ok=True)
IMAGES_DIR.mkdir(parents=True, exist_ok=True)


# ═══════════════════════════════════════════════════════════════
# 🎨 COLOR THEMES - Professional Design Palettes
# ═══════════════════════════════════════════════════════════════

COLOR_THEMES = {
    "professional": {
        "name": "Professional Blue",
        "primary": "#1E3A8A",      # Deep Blue
        "secondary": "#3B82F6",    # Bright Blue
        "accent": "#F59E0B",       # Amber
        "background": "#FFFFFF",   # White
        "text": "#1F2937",         # Dark Gray
        "light": "#E0F2FE",        # Light Blue
        "gradient_start": "#1E3A8A",
        "gradient_end": "#3B82F6"
    },
    "creative": {
        "name": "Creative Purple",
        "primary": "#7C3AED",      # Purple
        "secondary": "#EC4899",    # Pink
        "accent": "#10B981",       # Emerald
        "background": "#FFFFFF",
        "text": "#1F2937",
        "light": "#F3E8FF",
        "gradient_start": "#7C3AED",
        "gradient_end": "#EC4899"
    },
    "minimal": {
        "name": "Minimal Gray",
        "primary": "#374151",      # Gray
        "secondary": "#6B7280",    # Medium Gray
        "accent": "#3B82F6",       # Blue accent
        "background": "#FFFFFF",
        "text": "#111827",
        "light": "#F3F4F6",
        "gradient_start": "#374151",
        "gradient_end": "#6B7280"
    },
    "bold": {
        "name": "Bold Orange",
        "primary": "#DC2626",      # Red
        "secondary": "#F97316",    # Orange
        "accent": "#14B8A6",       # Teal
        "background": "#FFFFFF",
        "text": "#1F2937",
        "light": "#FEF3C7",
        "gradient_start": "#DC2626",
        "gradient_end": "#F97316"
    },
    "nature": {
        "name": "Nature Green",
        "primary": "#059669",      # Emerald
        "secondary": "#10B981",    # Green
        "accent": "#F59E0B",       # Amber
        "background": "#FFFFFF",
        "text": "#1F2937",
        "light": "#D1FAE5",
        "gradient_start": "#059669",
        "gradient_end": "#10B981"
    }
}


# ═══════════════════════════════════════════════════════════════
# 🚀 AI SLIDE GENERATOR CLASS
# ═══════════════════════════════════════════════════════════════

class AISlideGenerator:
    """
    AI Slide Generator with 7-Step Thinking Process
    สร้าง Presentation สวยงามแบบ Z.AI
    """
    
    def __init__(self, api_key: Optional[str] = None):
        """Initialize generator with GLM-4.7 API"""
        self.api_key = api_key or GLM_API_KEY
        self.base_url = "https://open.bigmodel.cn/api/paas/v4"
        self.model = "glm-4.7"
        
        # Image generation
        self.image_api_url = "https://open.bigmodel.cn/api/paas/v4/images/generations"
        self.cogview_model = "cogview-3-flash"
        
        # Azure DALL-E (fallback)
        self.azure_dalle_available = bool(AZURE_OPENAI_API_KEY and AZURE_DALLE_ENDPOINT)
        
        print(f"\n{'='*60}")
        print("🎯 AI Slide Generator (Z.AI Style)")
        print(f"{'='*60}")
        print(f"✅ Model: {self.model}")
        print(f"✅ PPTX Export: {'Ready' if PPTX_AVAILABLE else 'Not Available'}")
        print(f"✅ Image Generation: CogView-3 + DALL-E")
        print(f"{'='*60}\n")
    
    def _call_api(
        self,
        messages: List[Dict],
        max_tokens: int = 4096,
        temperature: float = 0.8
    ) -> str:
        """Call GLM API with retry logic - uses faster models"""
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }
        
        # Try faster models in order
        models_to_try = ["glm-4-flash", "glm-4-air", self.model]
        
        for model in models_to_try:
            payload = {
                "model": model,
                "messages": messages,
                "temperature": temperature,
                "max_tokens": max_tokens,
                "stream": False
            }
            
            try:
                print(f"   🔄 Trying {model}...")
                
                response = requests.post(
                    f"{self.base_url}/chat/completions",
                    headers=headers,
                    json=payload,
                    timeout=60  # Shorter timeout
                )
                
                if response.status_code != 200:
                    error_text = response.text[:200]
                    print(f"   ⚠️ {model} error: {error_text}")
                    continue
                
                data = response.json()
                
                if "choices" not in data or len(data["choices"]) == 0:
                    print(f"   ⚠️ {model} returned empty choices")
                    continue
                
                message = data["choices"][0].get("message", {})
                content = message.get("content", "")
                
                # Handle reasoning mode
                if not content and message.get("reasoning_content"):
                    reasoning = message.get("reasoning_content", "")
                    import re
                    json_match = re.search(r'\{.*\}', reasoning, re.DOTALL)
                    if json_match:
                        content = json_match.group(0)
                
                if content and len(content) > 10:
                    print(f"   ✅ {model} response ({len(content)} chars)")
                    return content
                    
            except requests.exceptions.Timeout:
                print(f"   ⏱️ {model} timeout, trying next...")
                continue
            except Exception as e:
                print(f"   ❌ {model} error: {e}")
                continue
        
        print("   ❌ All models failed")
        return ""
    
    def _parse_json(self, content: str) -> Dict:
        """Parse JSON from response with robust error handling"""
        if not content or len(content) < 10:
            return {}
        
        content = content.strip()
        
        # Remove markdown code blocks
        if "```json" in content:
            start = content.find("```json") + 7
            end = content.find("```", start)
            if end > start:
                content = content[start:end].strip()
        elif "```" in content:
            start = content.find("```") + 3
            end = content.find("```", start)
            if end > start:
                content = content[start:end].strip()
        
        # Find JSON object
        if not content.startswith("{"):
            start = content.find("{")
            end = content.rfind("}") + 1
            if start != -1 and end > start:
                content = content[start:end]
            else:
                return {}
        
        try:
            return json.loads(content)
        except json.JSONDecodeError:
            import re
            # Try to fix common issues
            fixed = re.sub(r',\s*}', '}', content)
            fixed = re.sub(r',\s*]', ']', fixed)
            try:
                return json.loads(fixed)
            except:
                return {}
    
    def generate_image(self, prompt: str, slide_id: str) -> Optional[str]:
        """Generate image using CogView-3"""
        try:
            enhanced_prompt = f"{prompt} Professional presentation visual, modern design, high quality, 16:9 aspect ratio"
            
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.api_key}"
            }
            
            payload = {
                "model": self.cogview_model,
                "prompt": enhanced_prompt[:500],
                "size": "1440x720",
                "quality": "standard",
                "n": 1
            }
            
            response = requests.post(
                self.image_api_url,
                headers=headers,
                json=payload,
                timeout=60
            )
            
            if response.status_code == 200:
                data = response.json()
                image_url = data.get("data", [{}])[0].get("url")
                
                if image_url:
                    # Save locally
                    local_path = self._save_image(image_url, slide_id)
                    print(f"   ✅ Image generated for slide {slide_id}")
                    return str(local_path) if local_path else None
            else:
                print(f"   ⚠️ CogView error: {response.status_code}")
                
        except Exception as e:
            print(f"   ⚠️ Image generation error: {e}")
        
        return None
    
    def _save_image(self, url: str, slide_id: str) -> Optional[Path]:
        """Save image locally"""
        try:
            response = requests.get(url, timeout=30)
            if response.status_code == 200:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"slide_{slide_id}_{timestamp}.png"
                filepath = IMAGES_DIR / filename
                with open(filepath, 'wb') as f:
                    f.write(response.content)
                return filepath
        except Exception as e:
            print(f"   ⚠️ Save error: {e}")
        return None
    
    def generate(
        self,
        topic: str,
        slide_count: int = 5,
        style: str = "professional",
        research_context: Optional[str] = None,
        generate_images: bool = False,
        language: str = "auto"
    ) -> Dict[str, Any]:
        """
        🚀 Generate Complete Presentation with 7-Step Thinking
        
        Args:
            topic: หัวข้อ Presentation
            slide_count: จำนวน slides
            style: สไตล์ (professional, creative, minimal, bold, nature)
            research_context: ข้อมูลจาก DeepResearch
            generate_images: สร้างรูปภาพหรือไม่
            language: ภาษา
            
        Returns:
            Dict with title, slides, pptx_path, metadata
        """
        
        print(f"\n{'='*60}")
        print(f"🎯 7-Step AI Slide Generation")
        print(f"{'='*60}")
        print(f"📌 Topic: {topic}")
        print(f"📊 Slides: {slide_count}")
        print(f"🎨 Style: {style}")
        print(f"🖼️ Images: {'Yes' if generate_images else 'No'}")
        print(f"{'='*60}\n")
        
        # Get color theme
        theme = COLOR_THEMES.get(style, COLOR_THEMES["professional"])
        
        # Step 1-6: Generate content with single comprehensive API call
        print("   📝 Step 1-6: Generating comprehensive content...")
        
        context_text = ""
        if research_context:
            context_text = f"\n\n### Research Context:\n{research_context[:3000]}"
        
        system_prompt = f"""คุณเป็นนักออกแบบ Presentation มืออาชีพระดับโลก
สร้าง slides ที่สวยงาม มีคุณภาพสูง น่าสนใจ

Design Theme: {theme['name']}
Primary Color: {theme['primary']}
Style: {style}

ตอบเป็น JSON เท่านั้น ไม่ต้องอธิบาย"""

        user_prompt = f"""สร้าง Presentation {slide_count} slides เกี่ยวกับ: {topic}
{context_text}

ตอบเป็น JSON format นี้:
{{
  "title": "หัวข้อหลัก",
  "subtitle": "คำบรรยายสั้นๆ",
  "slides": [
    {{
      "slide_number": 1,
      "type": "title",
      "title": "หัวข้อหลัก",
      "subtitle": "คำบรรยาย",
      "content": [],
      "speaker_notes": "ต้อนรับผู้ชม...",
      "design": {{
        "layout": "center",
        "has_image": false
      }}
    }},
    {{
      "slide_number": 2,
      "type": "content",
      "title": "หัวข้อ",
      "subtitle": "",
      "content": [
        "ประเด็นสำคัญข้อ 1 (10-15 คำ)",
        "ประเด็นสำคัญข้อ 2",
        "ประเด็นสำคัญข้อ 3"
      ],
      "speaker_notes": "อธิบายเพิ่มเติม...",
      "design": {{
        "layout": "two-column",
        "has_image": true,
        "image_prompt": "คำอธิบายรูปที่ต้องการ"
      }}
    }},
    {{
      "slide_number": 3,
      "type": "quote",
      "title": "",
      "quote": "คำพูดที่น่าจดจำ",
      "author": "ผู้พูด",
      "speaker_notes": "",
      "design": {{
        "layout": "center"
      }}
    }},
    {{
      "slide_number": 4,
      "type": "list",
      "title": "หัวข้อ",
      "items": ["รายการ 1", "รายการ 2", "รายการ 3"],
      "speaker_notes": "",
      "design": {{
        "layout": "grid"
      }}
    }},
    {{
      "slide_number": {slide_count},
      "type": "closing",
      "title": "สรุปและขอบคุณ",
      "content": ["ประเด็นสำคัญ", "Call to action"],
      "speaker_notes": "ขอบคุณผู้ฟัง"
    }}
  ],
  "metadata": {{
    "total_slides": {slide_count},
    "style": "{style}",
    "language": "auto"
  }}
}}

กฎ:
- เนื้อหาแต่ละข้อ 10-15 คำ
- มี 3-5 bullet points ต่อ slide
- ใช้ภาษาเดียวกับหัวข้อ
- สร้างเนื้อหาที่น่าสนใจ มีคุณค่า
- Slide แรกเป็น title, slide สุดท้ายเป็น closing/thank you"""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
        
        response = self._call_api(messages, max_tokens=6000, temperature=0.8)
        result = self._parse_json(response)
        
        if not result or "slides" not in result:
            print("   ⚠️ Using fallback content...")
            result = self._create_fallback_presentation(topic, slide_count, style)
        
        slides = result.get("slides", [])
        
        # Ensure slides have required fields
        for i, slide in enumerate(slides):
            slide["slide_number"] = i + 1
            if "type" not in slide:
                slide["type"] = "title" if i == 0 else "content"
            if "content" not in slide:
                slide["content"] = []
            if "speaker_notes" not in slide:
                slide["speaker_notes"] = ""
        
        # Step 7: Generate images if requested
        if generate_images:
            print("   🎨 Step 7: Generating images...")
            for slide in slides:
                design = slide.get("design", {})
                if design.get("has_image") and design.get("image_prompt"):
                    image_path = self.generate_image(
                        design["image_prompt"],
                        str(slide["slide_number"])
                    )
                    if image_path:
                        slide["image_path"] = image_path
        
        print(f"   ✅ Generated {len(slides)} slides successfully!")
        
        # Create PPTX
        pptx_path = None
        if PPTX_AVAILABLE:
            print("   💾 Creating PPTX file...")
            pptx_path = self._create_pptx(
                title=result.get("title", topic),
                slides=slides,
                theme=theme,
                topic=topic
            )
            print(f"   ✅ PPTX saved: {pptx_path}")
        
        return {
            "success": True,
            "title": result.get("title", topic),
            "subtitle": result.get("subtitle", ""),
            "slides": slides,
            "pptx_path": str(pptx_path) if pptx_path else None,
            "metadata": {
                "model": "GLM-4.7 (Z.AI 7-Step)",
                "style": style,
                "theme": theme["name"],
                "total_slides": len(slides),
                "has_images": generate_images,
                "generated_at": datetime.now().isoformat()
            }
        }
    
    def _create_fallback_presentation(
        self,
        topic: str,
        slide_count: int,
        style: str
    ) -> Dict[str, Any]:
        """Create fallback presentation when API fails"""
        slides = [
            {
                "slide_number": 1,
                "type": "title",
                "title": topic,
                "subtitle": "Professional Presentation",
                "content": [],
                "speaker_notes": "Welcome to this presentation"
            }
        ]
        
        for i in range(2, slide_count):
            slides.append({
                "slide_number": i,
                "type": "content",
                "title": f"Key Point {i-1}",
                "content": [
                    f"Important aspect of {topic}",
                    "Supporting information and details",
                    "Implementation strategies"
                ],
                "speaker_notes": f"Discuss key point {i-1}"
            })
        
        slides.append({
            "slide_number": slide_count,
            "type": "closing",
            "title": "Thank You",
            "content": ["Questions?", "Contact Information"],
            "speaker_notes": "Thank the audience"
        })
        
        return {
            "title": topic,
            "subtitle": "",
            "slides": slides
        }
    
    def _create_pptx(
        self,
        title: str,
        slides: List[Dict],
        theme: Dict[str, str],
        topic: str
    ) -> Optional[Path]:
        """Create beautiful PPTX file"""
        if not PPTX_AVAILABLE:
            return None
        
        try:
            prs = PPTXPresentation()
            prs.slide_width = Inches(13.333)  # 16:9
            prs.slide_height = Inches(7.5)
            
            for slide_data in slides:
                self._add_slide(prs, slide_data, theme)
            
            # Save file
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_topic = "".join(c for c in topic[:30] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
            filename = f"{timestamp}_{safe_topic}.pptx"
            filepath = EXPORTS_DIR / filename
            
            prs.save(str(filepath))
            return filepath
            
        except Exception as e:
            print(f"   ❌ PPTX creation error: {e}")
            return None
    
    def _add_slide(
        self,
        prs: PPTXPresentation,
        slide_data: Dict,
        theme: Dict[str, str]
    ):
        """Add a single slide to presentation with beautiful design"""
        try:
            layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(layout)
            
            slide_type = slide_data.get("type", "content")
            
            # Colors
            primary_color = self._hex_to_rgb(theme["primary"])
            text_color = self._hex_to_rgb(theme["text"])
            light_color = self._hex_to_rgb(theme["light"])
            accent_color = self._hex_to_rgb(theme["accent"])
            
            if slide_type == "title":
                self._create_title_slide(slide, slide_data, theme, primary_color, text_color)
            elif slide_type == "quote":
                self._create_quote_slide(slide, slide_data, theme, primary_color, text_color)
            elif slide_type == "closing":
                self._create_closing_slide(slide, slide_data, theme, primary_color, text_color)
            else:
                self._create_content_slide(slide, slide_data, theme, primary_color, text_color, light_color)
                
        except Exception as e:
            print(f"   ⚠️ Error adding slide: {e}")
    
    def _create_title_slide(self, slide, slide_data: Dict, theme: Dict, primary_color, text_color):
        """Create beautiful title slide"""
        # Gradient-style background strip
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = primary_color
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3),
            Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "Presentation")
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.7),
                Inches(11.333), Inches(0.8)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self._hex_to_rgb(theme["secondary"])
            p.alignment = PP_ALIGN.CENTER
        
        # Bottom accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.667), Inches(6.8),
            Inches(4), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(theme["accent"])
        line.line.fill.background()
    
    def _create_content_slide(self, slide, slide_data: Dict, theme: Dict, primary_color, text_color, light_color):
        """Create content slide with modern design"""
        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.15), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary_color
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4),
            Inches(12), Inches(0.9)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = primary_color
        
        # Title underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.35),
            Inches(3), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(theme["accent"])
        line.line.fill.background()
        
        # Content area
        content = slide_data.get("content", [])
        if isinstance(content, list):
            content_text = "\n".join([f"• {item}" for item in content])
        else:
            content_text = str(content)
        
        # Check if slide has image
        has_image = slide_data.get("image_path")
        
        if has_image:
            # Two column layout
            content_width = 6
            
            # Add image
            try:
                pic = slide.shapes.add_picture(
                    slide_data["image_path"],
                    Inches(7), Inches(1.5),
                    width=Inches(5.5), height=Inches(5)
                )
            except:
                content_width = 12
        else:
            content_width = 12
        
        if content_text:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8),
                Inches(content_width), Inches(5.2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            lines = content_text.split('\n')
            for idx, line in enumerate(lines):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = text_color
                p.space_before = Pt(12)
                p.space_after = Pt(8)
        
        # Slide number
        num_box = slide.shapes.add_textbox(
            Inches(12.5), Inches(7),
            Inches(0.5), Inches(0.3)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_data.get("slide_number", ""))
        p.font.size = Pt(12)
        p.font.color.rgb = self._hex_to_rgb(theme["secondary"])
        p.alignment = PP_ALIGN.RIGHT
    
    def _create_quote_slide(self, slide, slide_data: Dict, theme: Dict, primary_color, text_color):
        """Create quote slide"""
        # Background shape
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2),
            Inches(13.333), Inches(3.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._hex_to_rgb(theme["light"])
        bg.line.fill.background()
        
        # Quote mark
        quote_mark = slide.shapes.add_textbox(
            Inches(0.8), Inches(2),
            Inches(1), Inches(1)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(120)
        p.font.color.rgb = primary_color
        p.font.bold = True
        
        # Quote text
        quote = slide_data.get("quote", slide_data.get("title", ""))
        quote_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.8),
            Inches(10.333), Inches(2)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote
        p.font.size = Pt(32)
        p.font.italic = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        # Author
        author = slide_data.get("author", "")
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(1), Inches(5),
                Inches(11.333), Inches(0.5)
            )
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {author}"
            p.font.size = Pt(20)
            p.font.color.rgb = self._hex_to_rgb(theme["secondary"])
            p.alignment = PP_ALIGN.CENTER
    
    def _create_closing_slide(self, slide, slide_data: Dict, theme: Dict, primary_color, text_color):
        """Create closing/thank you slide"""
        # Full background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = primary_color
        bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "Thank You")
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Content (questions, contact info)
        content = slide_data.get("content", [])
        if content:
            if isinstance(content, list):
                content_text = "\n".join(content)
            else:
                content_text = str(content)
            
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2),
                Inches(11.333), Inches(2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for idx, line in enumerate(content_text.split('\n')):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
    
    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        """Convert hex color to RGB"""
        try:
            hex_color = hex_color.lstrip('#')
            if len(hex_color) == 6:
                r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                return RGBColor(r, g, b)
        except:
            pass
        return RGBColor(0, 0, 0)


# ═══════════════════════════════════════════════════════════════
# 🎯 PUBLIC FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def generate_ai_slides(
    topic: str,
    slide_count: int = 5,
    style: str = "professional",
    research_context: Optional[str] = None,
    generate_images: bool = False,
    user_id: Optional[str] = None
) -> Dict[str, Any]:
    """
    Generate AI Slides with 7-Step Process + PPTX Export
    
    Args:
        topic: หัวข้อ
        slide_count: จำนวน slides
        style: สไตล์ (professional, creative, minimal, bold, nature)
        research_context: ข้อมูลจาก research
        generate_images: สร้างรูปภาพ
        user_id: ID ผู้ใช้
        
    Returns:
        Dict with title, slides, pptx_path
    """
    generator = AISlideGenerator()
    result = generator.generate(
        topic=topic,
        slide_count=slide_count,
        style=style,
        research_context=research_context,
        generate_images=generate_images
    )
    
    # Save JSON to user folder if user_id provided
    if user_id and result.get("success"):
        user_dir = PROJECT_ROOT / "user_data" / user_id.replace("@", "_at_").replace(".", "_") / "presentations"
        user_dir.mkdir(parents=True, exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_topic = "".join(c for c in topic[:30] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
        json_path = user_dir / f"{timestamp}_{safe_topic}.json"
        
        # Copy PPTX to user folder too
        if result.get("pptx_path"):
            import shutil
            pptx_filename = Path(result["pptx_path"]).name
            user_pptx_path = user_dir / pptx_filename
            try:
                shutil.copy(result["pptx_path"], user_pptx_path)
                result["pptx_path"] = str(user_pptx_path)
            except:
                pass
        
        # Save JSON
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(result, f, indent=2, ensure_ascii=False)
        
        result["json_path"] = str(json_path)
    
    return result


# Module initialization
print("\n" + "="*60)
print("🎯 AI SLIDE GENERATOR (Z.AI STYLE)")
print("="*60)
print(f"✅ 7-Step Thinking Process Ready")
print(f"✅ Beautiful PPTX Export: {'Ready' if PPTX_AVAILABLE else 'Not Available'}")
print(f"✅ Exports: {EXPORTS_DIR}")
print("="*60 + "\n")
